import re
import datetime
import numpy as np
import pandas as pd
from io import BytesIO
from fpdf import FPDF
from utils.data_processing import fmt_currency, fmt_num
from utils.recommendations import get_patterns_text, get_recommendations_text

_EXPORT_COL_LABELS = {
    "creative_id": "Creative ID",
    "platform": "Platform",
    "format_concept": "Format / Concept",
    "length": "Length",
    "spend": "Spend ($)",
    "paid_starts": "Paid Starts",
    "trial_starts": "Trial Starts",
    "cpa": "Cost / Paid Start ($)",
    "cpt": "Cost / Trial Start ($)",
    "trial_to_paid_cvr": "Trial->Paid CVR (%)",
    "paid_per_1k": "Paid Starts / $1k",
    "trial_per_1k": "Trial Starts / $1k",
    "thumbstop_rate": "Thumbstop Rate (%)",
    "hold_6s": "6s Hold Rate (%)",
    "ctr": "CTR (%)",
    "goal_score": "Goal Score",
    "decision_label": "Decision",
}


def build_export_csv(df: pd.DataFrame) -> bytes:
    present = [c for c in _EXPORT_COL_LABELS if c in df.columns]
    out = df[present].copy().reset_index(drop=True)
    out.index = out.index + 1
    out.index.name = "Rank"
    for col in ["spend", "cpa", "cpt"]:
        if col in out.columns:
            out[col] = out[col].apply(lambda v: round(v, 2) if pd.notna(v) else v)
    for col in ["thumbstop_rate", "hold_6s", "ctr", "trial_to_paid_cvr"]:
        if col in out.columns:
            out[col] = out[col].apply(lambda v: round(v, 2) if pd.notna(v) else v)
    for col in ["paid_per_1k", "trial_per_1k", "goal_score"]:
        if col in out.columns:
            out[col] = out[col].apply(lambda v: round(v, 4) if pd.notna(v) else v)
    out = out.rename(columns=_EXPORT_COL_LABELS)
    return out.to_csv().encode("utf-8")


_UNICODE_REPLACE = str.maketrans({
    "\u2014": "-", "\u2013": "-", "\u2026": "...",
    "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
    "\u2022": "-", "\u00b7": "-", "\u2192": "->", "\u2190": "<-",
    "\u2264": "<=", "\u2265": ">=", "\u00a9": "(c)",
    "\u00ae": "(R)", "\u00b0": " deg",
})


def _strip_html(text: str) -> str:
    return re.sub(r"<[^>]+>", "", text).strip()


def _pdf_safe(text: str) -> str:
    text = _strip_html(text)
    text = text.translate(_UNICODE_REPLACE)
    return text.encode("latin-1", errors="replace").decode("latin-1")


def sanitize_pdf_text(text) -> str:
    if text is None:
        return ""
    text = str(text)
    replacements = {
        "\u2022": "-", "\u2013": "-", "\u2014": "-", "\u2015": "-", "\u2212": "-",
        "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
        "\u2026": "...", "\u2192": "->", "\u2190": "<-",
        "\u2264": "<=", "\u2265": ">=", "\u00a0": " ",
        "\u00ae": "(R)", "\u2122": "(TM)", "\u00a9": "(c)",
        "\u00b0": " deg", "\u00b7": "-",
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    encoded = text.encode("latin-1", "replace").decode("latin-1")
    result = []
    for orig_ch, enc_ch in zip(text, encoded):
        result.append("-" if enc_ch == "?" and orig_ch != "?" else enc_ch)
    if len(encoded) > len(text):
        result.extend(encoded[len(text):])
    return "".join(result)


def _pdf_cell(pdf, w, h, txt="", border=0, ln=0, align="", fill=False):
    pdf.cell(w, h, sanitize_pdf_text(txt), border=border, ln=ln, align=align, fill=fill)


def _pdf_multi_cell(pdf, w, h, txt="", border=0, align="", fill=False):
    pdf.multi_cell(w, h, sanitize_pdf_text(txt), border=border, align=align, fill=fill)


_PDF_DECISION_COLORS = {
    "Scale":        ((22, 101, 52),   (220, 252, 231)),
    "Keep Testing": ((146, 64, 14),   (254, 243, 199)),
    "Fix Funnel":   ((154, 52, 18),   (255, 237, 213)),
    "Cut":          ((153, 27, 27),   (254, 226, 226)),
    "Review":       ((120, 53, 15),   (254, 249, 195)),
}


class _PDF(FPDF):
    def __init__(self, goal: str, report_date: str):
        super().__init__(orientation="P", unit="mm", format="A4")
        self._goal = goal
        self._report_date = report_date
        self.set_margins(15, 15, 15)
        self.set_auto_page_break(auto=True, margin=15)

    def header(self):
        self.set_fill_color(30, 58, 138)
        self.rect(0, 0, 210, 18, "F")
        self.set_font("Helvetica", "B", 12)
        self.set_text_color(255, 255, 255)
        self.set_xy(15, 4)
        self.cell(120, 10, "Creative Performance Analyzer", ln=False)
        self.set_font("Helvetica", "", 9)
        self.set_xy(135, 4)
        self.cell(60, 5, sanitize_pdf_text(f"Goal: {self._goal}"), ln=False, align="R")
        self.set_xy(135, 9)
        self.cell(60, 5, sanitize_pdf_text(f"Generated: {self._report_date}"), ln=False, align="R")
        self.set_text_color(0, 0, 0)
        self.set_xy(15, 22)

    def footer(self):
        self.set_y(-12)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(150, 150, 150)
        self.cell(0, 5, sanitize_pdf_text(f"Page {self.page_no()} - Creative Performance Analyzer"), align="C")
        self.set_text_color(0, 0, 0)

    def section_title(self, title: str):
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(30, 58, 138)
        self.cell(0, 7, sanitize_pdf_text(title).upper(), ln=True)
        self.set_draw_color(30, 58, 138)
        self.set_line_width(0.4)
        self.line(15, self.get_y(), 195, self.get_y())
        self.ln(3)
        self.set_text_color(0, 0, 0)
        self.set_draw_color(0, 0, 0)
        self.set_line_width(0.2)


def build_pdf_report(df: pd.DataFrame, goal: str, cpa_target: float) -> bytes:
    today = datetime.date.today().strftime("%B %d, %Y")
    pdf = _PDF(goal=goal, report_date=today)
    pdf.add_page()

    total_spend = df["spend"].sum() if "spend" in df.columns else float("nan")
    total_trials = df["trial_starts"].sum() if "trial_starts" in df.columns else float("nan")
    total_paid = df["paid_starts"].sum() if "paid_starts" in df.columns else float("nan")
    blended_cpa = (total_spend / total_paid) if (pd.notna(total_spend) and total_paid) else float("nan")
    best_creative = str(df.iloc[0]["creative_id"]) if (not df.empty and "creative_id" in df.columns) else "-"
    scale_n = int((df["decision_label"] == "Scale").sum()) if "decision_label" in df.columns else 0
    cut_n = int((df["decision_label"] == "Cut").sum()) if "decision_label" in df.columns else 0

    cards = [
        ("Total Spend", fmt_currency(total_spend)),
        ("Trial Starts", fmt_num(total_trials)),
        ("Paid Starts", fmt_num(total_paid)),
        ("Blended CPA", fmt_currency(blended_cpa)),
        ("Best Creative", best_creative),
        ("Scale Candidates", str(scale_n)),
        ("Cut Candidates", str(cut_n)),
    ]

    pdf.section_title("Performance Summary")
    card_w = 180 / 7
    card_h = 14
    x0 = 15
    for i, (label, value) in enumerate(cards):
        x = x0 + i * card_w
        y = pdf.get_y()
        pdf.set_fill_color(241, 245, 249)
        pdf.set_draw_color(203, 213, 225)
        pdf.set_line_width(0.3)
        pdf.rect(x, y, card_w - 1, card_h, "FD")
        pdf.set_xy(x, y + 1)
        pdf.set_font("Helvetica", "", 6.5)
        pdf.set_text_color(100, 116, 139)
        _pdf_cell(pdf, card_w - 1, 4, label, align="C")
        pdf.set_xy(x, y + 5.5)
        pdf.set_font("Helvetica", "B", 9)
        pdf.set_text_color(15, 23, 42)
        _pdf_cell(pdf, card_w - 1, 6, value, align="C")
    pdf.ln(card_h + 4)
    pdf.set_text_color(0, 0, 0)

    pdf.section_title("Creative Rankings")
    tbl_cols = [
        ("Rank", 9, "C"), ("Creative ID", 30, "L"), ("Platform", 25, "L"),
        ("Spend ($)", 22, "R"), ("Paid Starts", 22, "R"), ("CPA ($)", 22, "R"),
        ("CTR (%)", 18, "R"), ("Decision", 32, "C"),
    ]
    col_keys = ["_rank", "creative_id", "platform", "spend", "paid_starts", "cpa", "ctr", "decision_label"]
    header_h, row_h = 7, 6
    ranked = df.reset_index(drop=True).copy()
    ranked.index = ranked.index + 1

    def _render_table_header():
        pdf.set_fill_color(30, 58, 138)
        pdf.set_text_color(255, 255, 255)
        pdf.set_font("Helvetica", "B", 7.5)
        for col_label, col_w, _ in tbl_cols:
            _pdf_cell(pdf, col_w, header_h, col_label, border=0, fill=True, align="C")
        pdf.ln(header_h)
        pdf.set_text_color(0, 0, 0)

    _render_table_header()
    for i, row in ranked.iterrows():
        if pdf.get_y() > 260:
            pdf.add_page()
            _render_table_header()
        decision = str(row.get("decision_label", "")) if "decision_label" in ranked.columns else ""
        txt_rgb, fill_rgb = _PDF_DECISION_COLORS.get(decision, ((50, 50, 50), (248, 250, 252)))
        fill_color = (248, 250, 252) if i % 2 == 0 else (255, 255, 255)
        for idx, (_, col_w, align) in enumerate(tbl_cols):
            key = col_keys[idx]
            if key == "_rank":
                val = str(i)
            elif key == "decision_label":
                val = decision
            elif key not in ranked.columns:
                val = "-"
            else:
                raw = row[key]
                if pd.isna(raw):
                    val = "-"
                elif key in ("spend", "cpa"):
                    val = f"${raw:,.0f}"
                elif key == "paid_starts":
                    val = f"{int(raw):,}"
                elif key == "ctr":
                    val = f"{raw:.1f}%"
                else:
                    val = str(raw)
            if key == "decision_label" and decision:
                pdf.set_fill_color(*fill_rgb)
                pdf.set_text_color(*txt_rgb)
            else:
                pdf.set_fill_color(*fill_color)
                pdf.set_text_color(30, 41, 59)
            pdf.set_font("Helvetica", "B" if key == "decision_label" else "", 7)
            _pdf_cell(pdf, col_w, row_h, sanitize_pdf_text(val)[:22], border=0, fill=True, align=align)
        pdf.ln(row_h)

    pdf.set_text_color(0, 0, 0)
    pdf.ln(4)

    pdf.section_title("Charts")
    pdf.set_font("Helvetica", "I", 8.5)
    pdf.set_text_color(100, 116, 139)
    _pdf_multi_cell(
        pdf, 0, 5.5,
        "Interactive charts are displayed in the app. Download individual chart data as CSV using the buttons in the Visuals section.",
        border=0,
    )
    pdf.set_text_color(0, 0, 0)
    pdf.ln(3)

    patterns = get_patterns_text(df)
    if patterns:
        if pdf.get_y() > 240:
            pdf.add_page()
        pdf.section_title("Pattern Analysis")
        for item in patterns:
            if pdf.get_y() > 270:
                pdf.add_page()
            pdf.set_font("Helvetica", "", 8.5)
            pdf.set_text_color(30, 41, 59)
            pdf.set_fill_color(239, 246, 255)
            pdf.set_x(15)
            pdf.set_draw_color(59, 130, 246)
            pdf.set_line_width(0.8)
            _pdf_multi_cell(pdf, 0, 5.5, f"  {sanitize_pdf_text(item)}", border="L", fill=True)
            pdf.ln(1.5)
        pdf.set_draw_color(0, 0, 0)
        pdf.set_line_width(0.2)
        pdf.ln(2)

    recs = get_recommendations_text(df)
    if recs:
        if pdf.get_y() > 240:
            pdf.add_page()
        pdf.section_title("Recommendations - What to Test Next")
        for rec in recs:
            if pdf.get_y() > 270:
                pdf.add_page()
            pdf.set_font("Helvetica", "", 8.5)
            pdf.set_text_color(30, 41, 59)
            pdf.set_fill_color(240, 253, 244)
            pdf.set_x(15)
            pdf.set_draw_color(34, 197, 94)
            pdf.set_line_width(0.8)
            _pdf_multi_cell(pdf, 0, 5.5, f"  >>  {sanitize_pdf_text(rec)}", border="L", fill=True)
            pdf.ln(1.5)

    return bytes(pdf.output())


def build_pptx_report(df: pd.DataFrame, goal: str, cpa_target: float) -> bytes:
    from pptx import Presentation as _Prs
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = _Prs()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _txt(slide, left, top, width, height, text, size=13, bold=False,
             color=(250, 250, 250), align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sanitize_pdf_text(str(text))
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        p.alignment = align

    def _bg(slide, color=(0x0E, 0x11, 0x17)):
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*color)
        bg.line.fill.background()

    s1 = prs.slides.add_slide(blank)
    _bg(s1)
    _txt(s1, 1.5, 2.0, 10, 1.4, "Creative Performance Report",
         size=36, bold=True, color=(0x4F, 0x8E, 0xF7), align=PP_ALIGN.CENTER)
    _txt(s1, 1.5, 3.6, 10, 0.7,
         f"Goal: {goal}  |  CPA Target: ${cpa_target:,.0f}  |  {datetime.date.today().strftime('%B %d, %Y')}",
         size=14, color=(0x8A, 0x9B, 0xC8), align=PP_ALIGN.CENTER)

    s2 = prs.slides.add_slide(blank)
    _bg(s2)
    _txt(s2, 0.5, 0.2, 12, 0.7, "Performance Summary", size=22, bold=True, color=(0x4F, 0x8E, 0xF7))
    metrics: list[tuple[str, str]] = []
    if "spend" in df.columns:
        metrics.append(("Total Spend", f"${df['spend'].sum(skipna=True):,.0f}"))
    if "paid_starts" in df.columns:
        tp = df["paid_starts"].sum(skipna=True)
        metrics.append(("Paid Starts", f"{int(tp):,}"))
        if "spend" in df.columns and tp > 0:
            metrics.append(("Blended CPA", f"${df['spend'].sum(skipna=True)/tp:,.2f}"))
    if not df.empty and "creative_id" in df.columns:
        metrics.append(("Top Creative", str(df.iloc[0]["creative_id"])))
    if "decision_label" in df.columns:
        metrics.append(("Scale", str((df["decision_label"] == "Scale").sum())))
        metrics.append(("Cut", str((df["decision_label"] == "Cut").sum())))
    cw = 1.8
    for i, (lbl, val) in enumerate(metrics[:6]):
        x = 0.5 + i * (cw + 0.22)
        box = s2.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(cw), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x27)
        box.line.color.rgb = RGBColor(0x2A, 0x33, 0x50)
        _txt(s2, x + 0.05, 1.28, cw - 0.1, 0.4, lbl, size=8, color=(0x8A, 0x9B, 0xC8), align=PP_ALIGN.CENTER)
        _txt(s2, x + 0.05, 1.65, cw - 0.1, 0.85, val, size=17, bold=True, color=(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    s3 = prs.slides.add_slide(blank)
    _bg(s3)
    _txt(s3, 0.5, 0.2, 12, 0.7, "Creative Rankings", size=22, bold=True, color=(0x4F, 0x8E, 0xF7))
    tbl_keys = [c for c in ["creative_id", "platform", "spend", "paid_starts", "cpa", "ctr", "decision_label"] if c in df.columns]
    tbl_hdrs = {
        "creative_id": "Creative", "platform": "Platform", "spend": "Spend ($)",
        "paid_starts": "Paid Starts", "cpa": "CPA ($)", "ctr": "CTR (%)", "decision_label": "Decision",
    }
    top10 = df.head(10).reset_index(drop=True)
    if tbl_keys and not top10.empty:
        tbl = s3.shapes.add_table(
            len(top10) + 1, len(tbl_keys), Inches(0.4), Inches(1.0), Inches(12.5), Inches(6.0)
        ).table
        for j, k in enumerate(tbl_keys):
            c = tbl.cell(0, j)
            c.text = tbl_hdrs.get(k, k)
            c.text_frame.paragraphs[0].font.bold = True
            c.text_frame.paragraphs[0].font.size = Pt(9)
            c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
        for i, (_, row) in enumerate(top10.iterrows()):
            for j, k in enumerate(tbl_keys):
                c = tbl.cell(i + 1, j)
                v = row.get(k, "")
                if k in ("spend", "cpa") and pd.notna(v):
                    c.text = f"${float(v):,.0f}"
                elif k == "ctr" and pd.notna(v):
                    c.text = f"{float(v):.1f}%"
                elif k == "paid_starts" and pd.notna(v):
                    c.text = f"{int(v):,}"
                else:
                    c.text = sanitize_pdf_text(str(v)) if pd.notna(v) else "-"
                c.text_frame.paragraphs[0].font.size = Pt(8)
                c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
                bg = (0xF1, 0xF5, 0xF9) if i % 2 == 0 else (0xFF, 0xFF, 0xFF)
                c.fill.solid()
                c.fill.fore_color.rgb = RGBColor(*bg)

    s4 = prs.slides.add_slide(blank)
    _bg(s4)
    _txt(s4, 0.5, 0.2, 12, 0.7, "Recommendations", size=22, bold=True, color=(0x4F, 0x8E, 0xF7))
    recs = get_recommendations_text(df)
    box = s4.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0))
    tf4 = box.text_frame
    tf4.word_wrap = True
    for i, rec in enumerate(recs[:14]):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.text = f"- {sanitize_pdf_text(rec)}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xC5, 0xCF, 0xDF)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
