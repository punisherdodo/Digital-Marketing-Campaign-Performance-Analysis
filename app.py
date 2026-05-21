import re
import math
import datetime
import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.subplots as psub
import plotly.graph_objects as go
from io import StringIO, BytesIO
from fpdf import FPDF
import json
import os
import html as _html

st.set_page_config(
    page_title="Creative Performance Analyzer",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ── Custom CSS for dark dashboard polish ────────────────────────────────────
st.markdown(
    """
    <style>
    /* Metric card look */
    div[data-testid="metric-container"] {
        background: #161B27;
        border: 1px solid #2A3350;
        border-radius: 10px;
        padding: 14px 18px 10px;
    }
    div[data-testid="metric-container"] label { color: #8A9BC8 !important; font-size: 0.78rem; }
    div[data-testid="metric-container"] div[data-testid="stMetricValue"] { font-size: 1.4rem; }

    /* Section headers */
    .section-header {
        font-size: 1.1rem;
        font-weight: 700;
        color: #4F8EF7;
        text-transform: uppercase;
        letter-spacing: 0.08em;
        margin: 0.5rem 0 0.25rem;
    }

    /* Decision label badges */
    .badge-scale    { background:#1a4731; color:#34d399; border:1px solid #34d399;
                      padding:2px 10px; border-radius:20px; font-size:0.78rem; font-weight:700; }
    .badge-keep     { background:#3b3515; color:#fbbf24; border:1px solid #fbbf24;
                      padding:2px 10px; border-radius:20px; font-size:0.78rem; font-weight:700; }
    .badge-review   { background:#2e2a10; color:#f59e0b; border:1px solid #f59e0b;
                      padding:2px 10px; border-radius:20px; font-size:0.78rem; font-weight:700; }
    .badge-fix      { background:#3b2200; color:#fb923c; border:1px solid #fb923c;
                      padding:2px 10px; border-radius:20px; font-size:0.78rem; font-weight:700; }
    .badge-cut      { background:#3b0a0a; color:#f87171; border:1px solid #f87171;
                      padding:2px 10px; border-radius:20px; font-size:0.78rem; font-weight:700; }

    /* Pattern analysis bullets */
    .insight-box {
        background: #161B27;
        border-left: 3px solid #4F8EF7;
        border-radius: 0 8px 8px 0;
        padding: 10px 16px;
        margin: 6px 0;
        font-size: 0.92rem;
        color: #FAFAFA;
    }
    .rec-box {
        background: #161B27;
        border-left: 3px solid #34d399;
        border-radius: 0 8px 8px 0;
        padding: 10px 16px;
        margin: 6px 0;
        font-size: 0.92rem;
        color: #FAFAFA;
    }

    /* Divider */
    hr { border-color: #2A3350 !important; margin: 1.2rem 0; }

    /* Hide default menu */
    #MainMenu { visibility: hidden; }
    footer { visibility: hidden; }
    </style>
    """,
    unsafe_allow_html=True,
)

# ── Sample data ──────────────────────────────────────────────────────────────
SAMPLE_CSV = """Creative ID,Platform,Format / Concept,Length,Spend,Thumbstop Rate,6s Hold Rate,CTR,Trial Starts,Paid Starts
C01,Meta,"UGC problem/solution: I used to leave every meeting with 4 tabs open...",28s,4800,36%,24%,1.7%,267,59
C02,Meta,"Static screenshots: AI notes in one tap",15s,3200,25%,N/A,0.9%,133,35
C03,Meta,"Founder talking head: Why we built LoopNote",35s,2700,18%,11%,0.6%,54,12
C04,TikTok,"Creator skit: POV your boss asks what was decided in the meeting",23s,5100,45%,31%,2.4%,300,42
C05,TikTok,"Screen-record demo: Turn a rambling voice memo into a task list",21s,4400,39%,27%,2.1%,338,71
C06,TikTok,"AI avatar explainer: Meet your AI productivity assistant",30s,3000,22%,15%,0.8%,64,9
C07,YouTube Shorts,"Before/after comparison: Before vs. after meeting notes",32s,4700,38%,22%,1.3%,235,63
C08,YouTube Shorts,"Tutorial: 3 ways students use LoopNote",45s,3900,33%,29%,1.1%,112,13
C09,Meta,"Creator testimonial: I imported 120 voice notes and found 18 tasks",26s,5500,41%,26%,2.0%,306,67
C10,TikTok,"Trend remix: Things I stopped doing after using LoopNote",18s,3800,50%,35%,1.9%,174,20
"""

# ── Column normalisation map ──────────────────────────────────────────────────
# Maps lowercased / stripped versions of common column names → canonical names
COLUMN_MAP = {
    "creative id": "creative_id",
    "creative_id": "creative_id",
    "creativeid": "creative_id",
    "id": "creative_id",
    "platform": "platform",
    "format / concept": "format_concept",
    "format/concept": "format_concept",
    "format concept": "format_concept",
    "format": "format_concept",
    "concept": "format_concept",
    "length": "length",
    "spend": "spend",
    "thumbstop rate": "thumbstop_rate",
    "thumbstop_rate": "thumbstop_rate",
    "thumbstop": "thumbstop_rate",
    "6s hold rate": "hold_6s",
    "6s_hold_rate": "hold_6s",
    "6s hold": "hold_6s",
    "hold rate": "hold_6s",
    "hold_rate": "hold_6s",
    "ctr": "ctr",
    "click through rate": "ctr",
    "click-through rate": "ctr",
    "trial starts": "trial_starts",
    "trial_starts": "trial_starts",
    "trials": "trial_starts",
    "paid starts": "paid_starts",
    "paid_starts": "paid_starts",
    "paid": "paid_starts",
}

CPA_TARGET = 90.0

# Default scoring weights (as percentages that must sum to 100)
CE_WEIGHTS_DEFAULT = {"thumbstop_rate": 40, "hold_6s": 30, "ctr": 30}
FFQ_WEIGHTS_DEFAULT = {"cpa": 30, "paid_starts": 25, "trial_to_paid_cvr": 20, "ctr": 15, "thumbstop_rate": 10}

SAVES_PATH = "saved_analyses.json"


# ── Save / load helpers ───────────────────────────────────────────────────────

def _load_saves_file() -> list[dict]:
    """Return the list of saved analyses from disk, or [] if none exist."""
    if not os.path.exists(SAVES_PATH):
        return []
    try:
        with open(SAVES_PATH, "r") as f:
            data = json.load(f)
        return data if isinstance(data, list) else []
    except Exception:
        return []


def _write_saves_file(saves: list[dict]) -> None:
    with open(SAVES_PATH, "w") as f:
        json.dump(saves, f, indent=2)


def list_saves() -> list[dict]:
    """Return saved analyses sorted newest-first."""
    saves = _load_saves_file()
    return sorted(saves, key=lambda s: s.get("saved_at", ""), reverse=True)


def save_analysis(name: str, df_raw: pd.DataFrame, goal: str, cpa_target: float) -> None:
    """Persist the current analysis under the given name."""
    saves = _load_saves_file()
    new_save = {
        "id": datetime.datetime.now().strftime("%Y%m%d_%H%M%S_%f"),
        "name": name.strip() or "Untitled",
        "saved_at": datetime.datetime.now().strftime("%Y-%m-%d %H:%M"),
        "csv_data": df_raw.to_csv(index=False),
        "goal": goal,
        "cpa_target": cpa_target,
    }
    saves.append(new_save)
    _write_saves_file(saves)


def delete_save(save_id: str) -> None:
    saves = _load_saves_file()
    saves = [s for s in saves if s.get("id") != save_id]
    _write_saves_file(saves)


def load_save(save_id: str) -> dict | None:
    saves = _load_saves_file()
    for s in saves:
        if s.get("id") == save_id:
            return s
    return None


def normalise_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Rename columns to canonical names using COLUMN_MAP."""
    rename = {}
    for col in df.columns:
        key = col.strip().lower()
        if key in COLUMN_MAP:
            rename[col] = COLUMN_MAP[key]
    return df.rename(columns=rename)


def clean_numeric(series: pd.Series) -> pd.Series:
    """Strip $, commas and convert to float. Coerce non-numeric → NaN."""
    cleaned = (
        series.astype(str)
        .str.replace(r"[$,]", "", regex=True)
        .str.strip()
        .replace({"N/A": np.nan, "n/a": np.nan, "NA": np.nan, "-": np.nan, "": np.nan})
    )
    return pd.to_numeric(cleaned, errors="coerce")


def clean_percent(series: pd.Series) -> pd.Series:
    """Convert percent strings like '36%' or '0.36' to float (0–100 scale)."""
    s = (
        series.astype(str)
        .str.replace(r"[$,]", "", regex=True)
        .str.strip()
        .replace({"N/A": np.nan, "n/a": np.nan, "NA": np.nan, "-": np.nan, "": np.nan})
    )
    result = []
    for v in s:
        try:
            if isinstance(v, str) and "%" in v:
                result.append(float(v.replace("%", "").strip()))
            else:
                val = float(v)
                # Only scale if clearly a proportion (0 < v <= 1)
                if 0 < val <= 1:
                    result.append(val * 100)
                else:
                    result.append(val)
        except (ValueError, TypeError):
            result.append(np.nan)
    return pd.Series(result, index=series.index)


def extract_seconds(series: pd.Series) -> pd.Series:
    """Convert length strings like '28s' or '28' to numeric seconds."""
    def parse(v):
        if pd.isna(v):
            return np.nan
        v = str(v).strip().lower()
        v = v.replace("s", "").replace("sec", "").strip()
        try:
            return float(v)
        except ValueError:
            return np.nan
    return series.apply(parse)


def load_and_clean(source) -> tuple[pd.DataFrame, list[str]]:
    """Load CSV/Excel/string, normalise columns, clean types. Returns (df, warnings)."""
    warnings = []
    if isinstance(source, str):
        df = pd.read_csv(StringIO(source))
    elif hasattr(source, "name") and source.name.endswith((".xlsx", ".xls")):
        df = pd.read_excel(source)
    else:
        try:
            df = pd.read_csv(source)
        except Exception:
            df = pd.read_excel(source)

    df = normalise_columns(df)

    # Ensure creative_id exists
    if "creative_id" not in df.columns:
        df.insert(0, "creative_id", [f"C{i+1:02d}" for i in range(len(df))])
        warnings.append("No 'Creative ID' column found — auto-generated IDs assigned.")

    # Clean Spend
    if "spend" in df.columns:
        df["spend"] = clean_numeric(df["spend"])

    # Clean percent columns
    for col in ["thumbstop_rate", "hold_6s", "ctr"]:
        if col in df.columns:
            df[col] = clean_percent(df[col])

    # Clean integer columns
    for col in ["trial_starts", "paid_starts"]:
        if col in df.columns:
            df[col] = pd.to_numeric(df[col], errors="coerce")

    # Clean length
    if "length" in df.columns:
        df["length_s"] = extract_seconds(df["length"])

    return df, warnings


def calculate_kpis(df: pd.DataFrame) -> tuple[pd.DataFrame, list[str]]:
    """Add KPI columns where source data exists. Returns (df, missing_warnings)."""
    missing = []

    def safe_div(num, den):
        result = np.where(den != 0, num / den, np.nan)
        return pd.Series(result, index=df.index)

    has = lambda c: c in df.columns and df[c].notna().any()

    if has("spend") and has("paid_starts"):
        df["cpa"] = safe_div(df["spend"], df["paid_starts"])
    else:
        missing.append("Cost per Paid Start (needs Spend + Paid Starts)")

    if has("spend") and has("trial_starts"):
        df["cpt"] = safe_div(df["spend"], df["trial_starts"])
    else:
        missing.append("Cost per Trial Start (needs Spend + Trial Starts)")

    if has("paid_starts") and has("trial_starts"):
        df["trial_to_paid_cvr"] = safe_div(df["paid_starts"], df["trial_starts"])
    else:
        missing.append("Trial→Paid CVR (needs Paid Starts + Trial Starts)")

    if has("paid_starts") and has("spend"):
        df["paid_per_1k"] = safe_div(df["paid_starts"], df["spend"]) * 1000
    else:
        missing.append("Paid Starts per $1k (needs Paid Starts + Spend)")

    if has("trial_starts") and has("spend"):
        df["trial_per_1k"] = safe_div(df["trial_starts"], df["spend"]) * 1000
    else:
        missing.append("Trial Starts per $1k (needs Trial Starts + Spend)")

    if has("ctr") and has("cpa"):
        df["ctr_efficiency"] = safe_div(df["ctr"], df["cpa"])
    else:
        missing.append("CTR Efficiency Score (needs CTR + Cost per Paid Start)")

    if has("thumbstop_rate") and has("cpa"):
        df["thumbstop_efficiency"] = safe_div(df["thumbstop_rate"], df["cpa"])
    else:
        missing.append("Thumbstop Efficiency (needs Thumbstop Rate + Cost per Paid Start)")

    if has("hold_6s") and has("cpa"):
        df["hold_efficiency"] = safe_div(df["hold_6s"], df["cpa"])
    else:
        missing.append("Hold Efficiency (needs 6s Hold Rate + Cost per Paid Start)")

    return df, missing


def _min_max(s: pd.Series, lower_is_better: bool = False) -> pd.Series:
    """Min-max normalise a series to [0, 1]. If lower_is_better, invert."""
    mn, mx = s.min(), s.max()
    if mx == mn:
        return pd.Series(0.5, index=s.index)
    norm = (s - mn) / (mx - mn)
    return (1 - norm) if lower_is_better else norm


def rank_by_goal(
    df: pd.DataFrame,
    goal: str,
    cpa_target: float = CPA_TARGET,
    ce_weights: dict | None = None,
    ffq_weights: dict | None = None,
) -> pd.DataFrame:
    """Return df sorted by the selected campaign goal, with a Goal Score column."""
    if ce_weights is None:
        ce_weights = CE_WEIGHTS_DEFAULT
    if ffq_weights is None:
        ffq_weights = FFQ_WEIGHTS_DEFAULT

    d = df.copy()
    if goal == "Paid Starts":
        if "paid_starts" in d.columns:
            d = d.sort_values("paid_starts", ascending=False)

    elif goal == "Trial Starts":
        if "trial_starts" in d.columns:
            d = d.sort_values("trial_starts", ascending=False)

    elif goal == "Efficient Paid Starts":
        if "cpa" in d.columns:
            d = d.sort_values(
                ["cpa", "paid_starts"], ascending=[True, False], na_position="last"
            )
            d["scale_candidate"] = d["cpa"] < cpa_target

    elif goal == "Efficient Trial Starts":
        if "cpt" in d.columns:
            d = d.sort_values(
                ["cpt", "trial_starts"], ascending=[True, False], na_position="last"
            )

    elif goal == "Creative Engagement":
        ce_col_order = ["thumbstop_rate", "hold_6s", "ctr"]
        cols = []
        weights = []
        for col in ce_col_order:
            if col in d.columns and col in ce_weights:
                cols.append(col)
                weights.append(ce_weights[col] / 100.0)
        if cols:
            total_w = sum(weights)
            score = sum(_min_max(d[c]) * w for c, w in zip(cols, weights))
            d["goal_score"] = score / total_w if total_w > 0 else 0
            d = d.sort_values("goal_score", ascending=False)

    elif goal == "Full Funnel Quality":
        ffq_col_order = [
            ("cpa", True),
            ("paid_starts", False),
            ("trial_to_paid_cvr", False),
            ("ctr", False),
            ("thumbstop_rate", False),
        ]
        components = [
            (c, ffq_weights.get(c, 0) / 100.0, inv)
            for c, inv in ffq_col_order
            if c in d.columns and ffq_weights.get(c, 0) > 0
        ]
        total_w = sum(w for c, w, _ in components)
        if total_w > 0:
            score = sum(
                _min_max(d[c], inv) * w
                for c, w, inv in components
            )
            d["goal_score"] = score / total_w
            d = d.sort_values("goal_score", ascending=False)

    d = d.reset_index(drop=True)
    d.index = d.index + 1
    return d


def assign_decision_labels(df: pd.DataFrame, cpa_target: float = CPA_TARGET) -> pd.DataFrame:
    """Add Decision_Label column."""
    d = df.copy()
    labels = []

    median_paid = d["paid_starts"].median() if "paid_starts" in d.columns else np.nan
    median_trial = d["trial_starts"].median() if "trial_starts" in d.columns else np.nan
    median_cvr = d["trial_to_paid_cvr"].median() if "trial_to_paid_cvr" in d.columns else np.nan
    median_thumbstop = d["thumbstop_rate"].median() if "thumbstop_rate" in d.columns else np.nan

    for _, row in d.iterrows():
        cpa = row.get("cpa", np.nan)
        paid = row.get("paid_starts", np.nan)
        trial = row.get("trial_starts", np.nan)
        cvr = row.get("trial_to_paid_cvr", np.nan)
        thumbstop = row.get("thumbstop_rate", np.nan)

        if not np.isnan(cpa) and not np.isnan(paid):
            if cpa < cpa_target and paid > median_paid:
                labels.append("Scale")
                continue
        if not np.isnan(thumbstop) and not np.isnan(paid):
            if thumbstop > median_thumbstop and paid < median_paid:
                labels.append("Keep Testing")
                continue
        if not np.isnan(trial) and not np.isnan(cvr):
            if trial > median_trial and cvr < median_cvr:
                labels.append("Fix Funnel")
                continue
        if not np.isnan(cpa) and not np.isnan(thumbstop):
            if cpa > cpa_target and thumbstop < median_thumbstop:
                labels.append("Cut")
                continue
        labels.append("Review")

    d["decision_label"] = labels
    return d


def badge_html(label: str) -> str:
    cls_map = {
        "Scale": "badge-scale",
        "Keep Testing": "badge-keep",
        "Review": "badge-review",
        "Fix Funnel": "badge-fix",
        "Cut": "badge-cut",
    }
    cls = cls_map.get(label, "badge-review")
    return f'<span class="{cls}">{label}</span>'


def fmt_currency(v):
    if pd.isna(v):
        return "—"
    return f"${v:,.0f}"


def fmt_pct(v):
    if pd.isna(v):
        return "—"
    return f"{v:.1f}%"


def fmt_num(v):
    if pd.isna(v):
        return "—"
    return f"{v:,.0f}"


def fmt_float(v, decimals=2):
    if pd.isna(v):
        return "—"
    return f"{v:,.{decimals}f}"


def length_bucket(s: float) -> str:
    if np.isnan(s):
        return "Unknown"
    if s < 20:
        return "Under 20s"
    elif s <= 30:
        return "20–30s"
    elif s <= 40:
        return "31–40s"
    else:
        return "Over 40s"


def render_summary_cards(df: pd.DataFrame, goal: str):
    total_spend = df["spend"].sum() if "spend" in df.columns else np.nan
    total_trials = df["trial_starts"].sum() if "trial_starts" in df.columns else np.nan
    total_paid = df["paid_starts"].sum() if "paid_starts" in df.columns else np.nan
    blended_cpa = (total_spend / total_paid) if (not np.isnan(total_spend) and total_paid) else np.nan

    best_creative = "—"
    if not df.empty and "creative_id" in df.columns:
        best_creative = str(df.iloc[0].get("creative_id", "—"))

    scale_n = (df["decision_label"] == "Scale").sum() if "decision_label" in df.columns else 0
    cut_n = (df["decision_label"] == "Cut").sum() if "decision_label" in df.columns else 0

    blended_cpa_str = (f"${blended_cpa:,.2f}" if not np.isnan(blended_cpa) else "—")

    row1 = st.columns(4)
    with row1[0]:
        st.metric("Spend", fmt_currency(total_spend))
    with row1[1]:
        st.metric("Trials", fmt_num(total_trials))
    with row1[2]:
        st.metric("Paid", fmt_num(total_paid))
    with row1[3]:
        st.metric("Blended CPA", blended_cpa_str)

    row2 = st.columns(3)
    with row2[0]:
        st.metric("Best Creative", best_creative)
    with row2[1]:
        st.metric("Scale", int(scale_n))
    with row2[2]:
        st.metric("Cut", int(cut_n))


_DECISION_LABEL_COLORS = {
    "Scale":        {"color": "#34d399", "background-color": "#1a4731"},
    "Keep Testing": {"color": "#fbbf24", "background-color": "#3b3515"},
    "Review":       {"color": "#f59e0b", "background-color": "#2e2a10"},
    "Fix Funnel":   {"color": "#fb923c", "background-color": "#3b2200"},
    "Cut":          {"color": "#f87171", "background-color": "#3b0a0a"},
}


def _style_decision_cell(val):
    styles = _DECISION_LABEL_COLORS.get(val, {})
    return "; ".join(f"{k}: {v}" for k, v in styles.items())


def render_ranking_table(df: pd.DataFrame):
    display_cols_map = {
        "creative_id": "Creative ID",
        "platform": "Platform",
        "format_concept": "Format / Concept",
        "spend": "Spend ($)",
        "paid_starts": "Paid Starts",
        "trial_starts": "Trial Starts",
        "cpa": "Cost / Paid Start ($)",
        "cpt": "Cost / Trial Start ($)",
        "trial_to_paid_cvr": "Trial→Paid CVR (%)",
        "thumbstop_rate": "Thumbstop Rate (%)",
        "hold_6s": "6s Hold Rate (%)",
        "ctr": "CTR (%)",
        "decision_label": "Decision",
    }

    present = [c for c in display_cols_map if c in df.columns]
    sub = df[present].copy().reset_index(drop=True)
    sub.index = sub.index + 1
    sub.index.name = "Rank"

    # Format display values
    for col in ["spend", "cpa", "cpt"]:
        if col in sub.columns:
            sub[col] = sub[col].apply(lambda v: round(v, 2) if pd.notna(v) else v)
    for col in ["thumbstop_rate", "hold_6s", "ctr", "trial_to_paid_cvr"]:
        if col in sub.columns:
            sub[col] = sub[col].apply(lambda v: round(v, 2) if pd.notna(v) else v)

    sub = sub.rename(columns=display_cols_map)
    display_label = display_cols_map.get("decision_label", "Decision")

    styler = sub.style
    if display_label in sub.columns:
        styler = styler.map(_style_decision_cell, subset=[display_label])

    st.markdown(
        "<p style='font-size:0.8rem;color:#8A9BC8;margin-bottom:4px;'>"
        "Click any column header to sort. "
        "<span style='color:#34d399'>■</span> Scale &nbsp;"
        "<span style='color:#fbbf24'>■</span> Keep Testing &nbsp;"
        "<span style='color:#f59e0b'>■</span> Review &nbsp;"
        "<span style='color:#fb923c'>■</span> Fix Funnel &nbsp;"
        "<span style='color:#f87171'>■</span> Cut"
        "</p>",
        unsafe_allow_html=True,
    )
    st.dataframe(styler, use_container_width=True, height=400)


def render_charts(df: pd.DataFrame, cpa_target: float = CPA_TARGET):
    chart_bg = "#0E1117"
    grid_color = "#1E2A45"
    font_color = "#FAFAFA"

    chart_cols = st.columns(2)

    # 1. CPA bar chart
    with chart_cols[0]:
        st.markdown("<div class='section-header'>Cost per Paid Start by Creative</div>", unsafe_allow_html=True)
        if "cpa" in df.columns and "creative_id" in df.columns:
            d = df[["creative_id", "cpa"]].dropna().sort_values("cpa")
            fig = px.bar(
                d, x="creative_id", y="cpa",
                labels={"creative_id": "Creative", "cpa": "CPA ($)"},
                color="cpa",
                color_continuous_scale=[[0, "#34d399"], [0.5, "#fbbf24"], [1, "#f87171"]],
            )
            fig.add_hline(
                y=cpa_target, line_dash="dash", line_color="#4F8EF7",
                annotation_text=f"${cpa_target:,.0f} Target", annotation_font_color="#4F8EF7",
            )
            fig.update_layout(
                paper_bgcolor=chart_bg, plot_bgcolor=chart_bg,
                font_color=font_color, showlegend=False,
                coloraxis_showscale=False,
                margin=dict(t=20, b=40, l=50, r=20),
                xaxis=dict(gridcolor=grid_color, tickangle=-30),
                yaxis=dict(gridcolor=grid_color),
            )
            st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("Need Spend + Paid Starts for this chart.")

    # 2. Paid Starts bar
    with chart_cols[1]:
        st.markdown("<div class='section-header'>Paid Starts by Creative</div>", unsafe_allow_html=True)
        if "paid_starts" in df.columns and "creative_id" in df.columns:
            d = df[["creative_id", "paid_starts", "platform"]].dropna(subset=["paid_starts"]).sort_values("paid_starts", ascending=False)
            color_map = {"Meta": "#4F8EF7", "TikTok": "#34d399", "YouTube Shorts": "#fb923c"}
            fig = px.bar(
                d, x="creative_id", y="paid_starts",
                labels={"creative_id": "Creative", "paid_starts": "Paid Starts"},
                color="platform" if "platform" in d.columns else None,
                color_discrete_map=color_map,
            )
            fig.update_layout(
                paper_bgcolor=chart_bg, plot_bgcolor=chart_bg,
                font_color=font_color, legend_title_text="Platform",
                margin=dict(t=20, b=40, l=50, r=20),
                xaxis=dict(gridcolor=grid_color, tickangle=-30),
                yaxis=dict(gridcolor=grid_color),
            )
            st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("Need Paid Starts for this chart.")

    chart_cols2 = st.columns(2)

    # 3. Scatter CTR vs CPA
    with chart_cols2[0]:
        st.markdown("<div class='section-header'>CTR vs Cost per Paid Start</div>", unsafe_allow_html=True)
        if "ctr" in df.columns and "cpa" in df.columns:
            d = df[["creative_id", "ctr", "cpa", "spend", "platform", "decision_label"]].dropna(subset=["ctr", "cpa"])
            size_col = d["spend"].fillna(d["spend"].median()) if "spend" in d.columns else None
            color_map = {"Meta": "#4F8EF7", "TikTok": "#34d399", "YouTube Shorts": "#fb923c"}
            fig = px.scatter(
                d, x="cpa", y="ctr",
                size="spend" if size_col is not None else None,
                color="platform" if "platform" in d.columns else None,
                color_discrete_map=color_map,
                text="creative_id",
                labels={"cpa": "CPA ($)", "ctr": "CTR (%)", "spend": "Spend"},
                hover_data=["creative_id", "decision_label"],
            )
            fig.add_vline(
                x=cpa_target, line_dash="dash", line_color="#4F8EF7",
                annotation_text=f"${cpa_target:,.0f} CPA target", annotation_font_color="#4F8EF7",
            )
            fig.update_traces(textposition="top center", textfont_size=9)
            fig.update_layout(
                paper_bgcolor=chart_bg, plot_bgcolor=chart_bg,
                font_color=font_color, legend_title_text="Platform",
                margin=dict(t=20, b=40, l=50, r=20),
                xaxis=dict(gridcolor=grid_color),
                yaxis=dict(gridcolor=grid_color),
            )
            st.plotly_chart(fig, use_container_width=True)
        else:
            st.info("Need CTR + Spend + Paid Starts for this chart.")

    # 4. Platform grouped
    with chart_cols2[1]:
        st.markdown("<div class='section-header'>Platform Performance Overview</div>", unsafe_allow_html=True)
        if "platform" in df.columns:
            agg = {}
            if "cpa" in df.columns:
                agg["Avg CPA"] = ("cpa", "mean")
            if "paid_starts" in df.columns:
                agg["Total Paid Starts"] = ("paid_starts", "sum")
            if "ctr" in df.columns:
                agg["Avg CTR"] = ("ctr", "mean")
            if "thumbstop_rate" in df.columns:
                agg["Avg Thumbstop"] = ("thumbstop_rate", "mean")
            if agg:
                plat = df.groupby("platform").agg(**agg).reset_index()
                plat_melt = plat.melt(id_vars="platform", var_name="Metric", value_name="Value")
                fig = px.bar(
                    plat_melt, x="platform", y="Value",
                    color="Metric", barmode="group",
                    labels={"platform": "Platform", "Value": "Value"},
                    color_discrete_sequence=["#4F8EF7", "#34d399", "#fbbf24", "#fb923c"],
                )
                fig.update_layout(
                    paper_bgcolor=chart_bg, plot_bgcolor=chart_bg,
                    font_color=font_color, legend_title_text="Metric",
                    margin=dict(t=20, b=40, l=50, r=20),
                    xaxis=dict(gridcolor=grid_color),
                    yaxis=dict(gridcolor=grid_color),
                )
                st.plotly_chart(fig, use_container_width=True)
            else:
                st.info("Need Platform + at least one metric for this chart.")
        else:
            st.info("No Platform column found.")


def render_patterns(df: pd.DataFrame):
    insights = []

    if "platform" in df.columns and "cpa" in df.columns:
        plat_cpa = df.groupby("platform")["cpa"].mean().dropna()
        if not plat_cpa.empty:
            best = plat_cpa.idxmin()
            worst = plat_cpa.idxmax()
            insights.append(f"<b>{best}</b> has the lowest avg CPA at <b>{fmt_currency(plat_cpa[best])}</b> — best platform for efficient paid starts.")
            if best != worst:
                insights.append(f"<b>{worst}</b> has the highest avg CPA at <b>{fmt_currency(plat_cpa[worst])}</b> — review spend allocation.")

    if "platform" in df.columns and "paid_starts" in df.columns:
        plat_paid = df.groupby("platform")["paid_starts"].sum().dropna()
        if not plat_paid.empty:
            best = plat_paid.idxmax()
            insights.append(f"<b>{best}</b> drives the most paid starts in total (<b>{fmt_num(plat_paid[best])}</b>).")

    if "format_concept" in df.columns and "cpa" in df.columns:
        concept_cpa = df.groupby("format_concept")["cpa"].mean().dropna()
        if not concept_cpa.empty:
            best_concept = concept_cpa.idxmin()
            short = best_concept[:60] + "…" if len(best_concept) > 60 else best_concept
            insights.append(f"Best concept by CPA: <b>\"{short}\"</b> at <b>{fmt_currency(concept_cpa[best_concept])}</b>.")

    if "platform" in df.columns and "cpa" in df.columns:
        plat_cpa = df.groupby("platform")["cpa"].mean().dropna()
        rows = " | ".join(f"{p}: {fmt_currency(v)}" for p, v in plat_cpa.items())
        insights.append(f"Average CPA by platform — {rows}")

    if "length_s" in df.columns and "cpa" in df.columns:
        d = df.dropna(subset=["length_s", "cpa"]).copy()
        d["bucket"] = d["length_s"].apply(length_bucket)
        bucket_cpa = d.groupby("bucket")["cpa"].mean().dropna()
        ORDER = ["Under 20s", "20–30s", "31–40s", "Over 40s"]
        bucket_cpa = bucket_cpa.reindex([b for b in ORDER if b in bucket_cpa.index])
        if not bucket_cpa.empty:
            rows = " | ".join(f"{b}: {fmt_currency(v)}" for b, v in bucket_cpa.items())
            best_b = bucket_cpa.idxmin()
            insights.append(f"Average CPA by length bucket — {rows}")
            insights.append(f"<b>{best_b}</b> creatives have the lowest average CPA.")

    if insights:
        for ins in insights:
            st.markdown(f"<div class='insight-box'>{ins}</div>", unsafe_allow_html=True)
    else:
        st.info("Upload data with Platform, Spend, and Paid Starts to see pattern analysis.")


def render_recommendations(df: pd.DataFrame):
    recs = []

    if "platform" in df.columns and "thumbstop_rate" in df.columns and "cpa" in df.columns:
        plat = df.groupby("platform").agg(
            thumbstop=("thumbstop_rate", "mean"),
            cpa=("cpa", "mean"),
            paid=("paid_starts", "sum") if "paid_starts" in df.columns else ("cpa", "count"),
        ).dropna()

        has_tiktok = "TikTok" in plat.index
        has_meta = "Meta" in plat.index

        if has_tiktok and has_meta:
            if plat.loc["TikTok", "thumbstop"] > plat.loc["Meta", "thumbstop"] and plat.loc["TikTok", "cpa"] > plat.loc["Meta", "cpa"]:
                recs.append("TikTok hooks are stopping scrolls but not converting. Test stronger product proof, clearer CTA, or a harder offer in the last 5 seconds of TikTok creatives.")
            if plat.loc["Meta", "thumbstop"] < plat.loc["TikTok", "thumbstop"] and plat.loc["Meta", "cpa"] < plat.loc["TikTok", "cpa"]:
                recs.append("Meta converts better despite lower thumbstop. Test more scroll-stopping hooks on Meta while keeping the same proven offer and CTA structure.")

    if "ctr" in df.columns and "cpa" in df.columns:
        median_ctr = df["ctr"].median()
        median_cpa = df["cpa"].median()
        high_ctr_high_cpa = df[(df["ctr"] > median_ctr) & (df["cpa"] > median_cpa)]
        if not high_ctr_high_cpa.empty:
            ids = ", ".join(high_ctr_high_cpa["creative_id"].astype(str).tolist())
            recs.append(f"Creatives <b>{ids}</b> have strong CTR but high CPA — the drop-off is likely in the landing page or trial-to-paid flow. Test a more direct offer page or reduce friction in the trial sign-up.")

    if "thumbstop_rate" in df.columns and "ctr" in df.columns:
        median_thumb = df["thumbstop_rate"].median()
        median_ctr = df["ctr"].median()
        high_thumb_low_ctr = df[(df["thumbstop_rate"] > median_thumb) & (df["ctr"] < median_ctr)]
        if not high_thumb_low_ctr.empty:
            ids = ", ".join(high_thumb_low_ctr["creative_id"].astype(str).tolist())
            recs.append(f"Creatives <b>{ids}</b> stop the scroll but don't earn the click. Test a clearer, more specific value proposition in the first 3 seconds — viewers are hooked but not yet convinced.")

    if "format_concept" in df.columns and "cpa" in df.columns:
        concept_cpa = df.groupby("format_concept")["cpa"].mean().dropna()
        if not concept_cpa.empty and len(concept_cpa) > 1:
            best_concept = concept_cpa.idxmin()
            short = best_concept[:60] + "…" if len(best_concept) > 60 else best_concept
            pct_better = ((concept_cpa.mean() - concept_cpa.min()) / concept_cpa.mean() * 100)
            recs.append(f"Concept \"{short}\" has the lowest CPA — <b>{pct_better:.0f}% below average</b>. Remix this format across other platforms and creative lengths to expand its reach.")

    if "platform" in df.columns and "cpa" in df.columns:
        plat_cpa = df.groupby("platform")["cpa"].mean().dropna()
        if len(plat_cpa) > 1:
            worst_plat = plat_cpa.idxmax()
            best_plat = plat_cpa.idxmin()
            recs.append(f"Consider shifting budget from <b>{worst_plat}</b> (avg CPA {fmt_currency(plat_cpa[worst_plat])}) toward <b>{best_plat}</b> (avg CPA {fmt_currency(plat_cpa[best_plat])}) until {worst_plat} creative improves.")

    if recs:
        for rec in recs:
            st.markdown(f"<div class='rec-box'>💡 {rec}</div>", unsafe_allow_html=True)
    else:
        st.info("Upload data to generate test recommendations.")


# ══════════════════════════════════════════════════════════════════════════════
# EXPORT HELPERS
# ══════════════════════════════════════════════════════════════════════════════

_EXPORT_COL_LABELS = {
    "creative_id": "Creative ID",
    "platform": "Platform",
    "format_concept": "Format / Concept",
    "length": "Length",
    "spend": "Spend ($)",
    "paid_starts": "Paid Starts",
    "trial_starts": "Trial Starts",
    "cpa": "Cost / Paid Start ($)",
    "cpt": "Cost / Trial Start ($)",
    "trial_to_paid_cvr": "Trial→Paid CVR (%)",
    "paid_per_1k": "Paid Starts / $1k",
    "trial_per_1k": "Trial Starts / $1k",
    "thumbstop_rate": "Thumbstop Rate (%)",
    "hold_6s": "6s Hold Rate (%)",
    "ctr": "CTR (%)",
    "goal_score": "Goal Score",
    "decision_label": "Decision",
}


def build_export_csv(df: pd.DataFrame) -> bytes:
    """Return a UTF-8 CSV byte string of the ranked dataframe, formatted for readability."""
    present = [c for c in _EXPORT_COL_LABELS if c in df.columns]
    out = df[present].copy().reset_index(drop=True)
    out.index = out.index + 1
    out.index.name = "Rank"

    for col in ["spend", "cpa", "cpt"]:
        if col in out.columns:
            out[col] = out[col].apply(lambda v: round(v, 2) if pd.notna(v) else v)
    for col in ["thumbstop_rate", "hold_6s", "ctr", "trial_to_paid_cvr"]:
        if col in out.columns:
            out[col] = out[col].apply(lambda v: round(v, 2) if pd.notna(v) else v)
    for col in ["paid_per_1k", "trial_per_1k", "goal_score"]:
        if col in out.columns:
            out[col] = out[col].apply(lambda v: round(v, 4) if pd.notna(v) else v)

    out = out.rename(columns=_EXPORT_COL_LABELS)
    return out.to_csv().encode("utf-8")


def build_charts_png(df: pd.DataFrame) -> bytes | None:
    """Build a 2×2 Plotly subplot image and return PNG bytes. Returns None if data is insufficient."""
    chart_bg = "#0E1117"
    grid_color = "#1E2A45"
    font_color = "#FAFAFA"
    color_map = {"Meta": "#4F8EF7", "TikTok": "#34d399", "YouTube Shorts": "#fb923c"}

    has_cpa = "cpa" in df.columns and "creative_id" in df.columns
    has_paid = "paid_starts" in df.columns and "creative_id" in df.columns
    has_scatter = "ctr" in df.columns and "cpa" in df.columns
    has_platform = "platform" in df.columns

    if not any([has_cpa, has_paid, has_scatter, has_platform]):
        return None

    fig = psub.make_subplots(
        rows=2, cols=2,
        subplot_titles=[
            "Cost per Paid Start by Creative",
            "Paid Starts by Creative",
            "CTR vs Cost per Paid Start",
            "Platform Performance Overview",
        ],
        vertical_spacing=0.18,
        horizontal_spacing=0.10,
    )

    # 1 — CPA bar
    if has_cpa:
        d = df[["creative_id", "cpa"]].dropna().sort_values("cpa")
        for _, row in d.iterrows():
            pct = (row["cpa"] - d["cpa"].min()) / max(d["cpa"].max() - d["cpa"].min(), 1e-9)
            r = int(52 + pct * (248 - 52))
            g = int(211 - pct * (211 - 113))
            b = int(153 - pct * (153 - 113))
            fig.add_trace(
                go.Bar(x=[row["creative_id"]], y=[row["cpa"]], marker_color=f"rgb({r},{g},{b})", showlegend=False),
                row=1, col=1,
            )
        fig.add_hline(y=CPA_TARGET, line_dash="dash", line_color="#4F8EF7", row=1, col=1)

    # 2 — Paid Starts bar
    if has_paid:
        d2 = df[["creative_id", "paid_starts", "platform"]].dropna(subset=["paid_starts"]).sort_values("paid_starts", ascending=False)
        seen_platforms = set()
        for _, row in d2.iterrows():
            plat = row.get("platform", "Other")
            color = color_map.get(plat, "#8A9BC8")
            show = plat not in seen_platforms
            seen_platforms.add(plat)
            fig.add_trace(
                go.Bar(x=[row["creative_id"]], y=[row["paid_starts"]], marker_color=color,
                       name=plat, legendgroup=plat, showlegend=show),
                row=1, col=2,
            )

    # 3 — Scatter CTR vs CPA
    if has_scatter:
        d3 = df[["creative_id", "ctr", "cpa", "spend", "platform"]].dropna(subset=["ctr", "cpa"])
        spend_vals = d3["spend"].fillna(d3["spend"].median()) if "spend" in d3.columns else pd.Series([10] * len(d3))
        size_norm = ((spend_vals - spend_vals.min()) / max(spend_vals.max() - spend_vals.min(), 1)) * 20 + 8
        for i, (_, row) in enumerate(d3.iterrows()):
            plat = row.get("platform", "Other")
            color = color_map.get(plat, "#8A9BC8")
            fig.add_trace(
                go.Scatter(
                    x=[row["cpa"]], y=[row["ctr"]],
                    mode="markers+text",
                    text=[row["creative_id"]],
                    textposition="top center",
                    marker=dict(size=float(size_norm.iloc[i]), color=color, opacity=0.85),
                    showlegend=False,
                ),
                row=2, col=1,
            )
        fig.add_vline(x=CPA_TARGET, line_dash="dash", line_color="#4F8EF7", row=2, col=1)

    # 4 — Platform grouped
    if has_platform:
        agg = {}
        if "cpa" in df.columns:
            agg["Avg CPA"] = ("cpa", "mean")
        if "paid_starts" in df.columns:
            agg["Total Paid Starts"] = ("paid_starts", "sum")
        if "ctr" in df.columns:
            agg["Avg CTR"] = ("ctr", "mean")
        if "thumbstop_rate" in df.columns:
            agg["Avg Thumbstop"] = ("thumbstop_rate", "mean")
        if agg:
            plat = df.groupby("platform").agg(**agg).reset_index()
            metric_colors = ["#4F8EF7", "#34d399", "#fbbf24", "#fb923c"]
            for mi, metric in enumerate([k for k in agg]):
                fig.add_trace(
                    go.Bar(
                        x=plat["platform"], y=plat[metric],
                        name=metric, marker_color=metric_colors[mi % len(metric_colors)],
                        legendgroup=f"metric_{metric}", showlegend=True,
                    ),
                    row=2, col=2,
                )
            fig.update_layout(barmode="group")

    fig.update_layout(
        paper_bgcolor=chart_bg,
        plot_bgcolor=chart_bg,
        font=dict(color=font_color, size=11),
        height=900,
        width=1400,
        title=dict(text="Creative Performance Summary", font=dict(size=16, color="#4F8EF7"), x=0.5),
        margin=dict(t=60, b=40, l=60, r=60),
    )
    for axis in fig.layout:
        if axis.startswith("xaxis") or axis.startswith("yaxis"):
            fig.layout[axis].update(gridcolor=grid_color)

    try:
        return fig.to_image(format="png", scale=1.5)
    except Exception as e:
        return e


def _strip_html(text: str) -> str:
    return re.sub(r"<[^>]+>", "", text).strip()


_UNICODE_REPLACE = str.maketrans({
    "\u2014": "-",   # em dash
    "\u2013": "-",   # en dash
    "\u2026": "...", # ellipsis
    "\u2018": "'",   # left single quote
    "\u2019": "'",   # right single quote
    "\u201c": '"',   # left double quote
    "\u201d": '"',   # right double quote
    "\u2022": "-",   # bullet -> hyphen
    "\u00b7": "-",   # middle dot -> hyphen
    "\u2192": "->",  # right arrow (used in Trial->Paid CVR)
    "\u2190": "<-",  # left arrow
    "\u2264": "<=",  # less-than-or-equal
    "\u2265": ">=",  # greater-than-or-equal
    "\u00a9": "(c)", # copyright
    "\u00ae": "(R)", # registered
    "\u00b0": " deg",# degree
})


def _pdf_safe(text: str) -> str:
    """Strip HTML, replace non-Latin-1 chars with ASCII equivalents."""
    text = _strip_html(text)
    text = text.translate(_UNICODE_REPLACE)
    return text.encode("latin-1", errors="replace").decode("latin-1")


def sanitize_pdf_text(text) -> str:
    """Comprehensive sanitiser: replace every non-Helvetica character with a
    safe ASCII equivalent, then encode/decode to catch anything remaining."""
    if text is None:
        return ""
    text = str(text)
    replacements = {
        "\u2022": "-",    # bullet
        "\u2013": "-",    # en dash
        "\u2014": "-",    # em dash
        "\u2015": "-",    # horizontal bar
        "\u2212": "-",    # minus sign
        "\u2018": "'",    # left single quote
        "\u2019": "'",    # right single quote
        "\u201c": '"',    # left double quote
        "\u201d": '"',    # right double quote
        "\u2026": "...",  # ellipsis
        "\u2192": "->",   # right arrow
        "\u2190": "<-",   # left arrow
        "\u2264": "<=",   # less-than-or-equal
        "\u2265": ">=",   # greater-than-or-equal
        "\u00a0": " ",    # non-breaking space
        "\u00ae": "(R)",  # registered trademark
        "\u2122": "(TM)", # trademark
        "\u00a9": "(c)",  # copyright
        "\u00b0": " deg", # degree sign
        "\u00b7": "-",    # middle dot
    }
    for old, new in replacements.items():
        text = text.replace(old, new)
    # Final pass: encode to latin-1, replacing any remaining unsupported chars
    # with "?" then convert those replacement markers to hyphens.
    encoded = text.encode("latin-1", "replace").decode("latin-1")
    # Only replace '?' chars that were substituted (not original question marks)
    # by comparing char-by-char with the original encoded result.
    result = []
    for orig_ch, enc_ch in zip(text, encoded):
        if enc_ch == "?" and orig_ch != "?":
            result.append("-")
        else:
            result.append(enc_ch)
    # Handle length differences (substitutions that expand/contract)
    if len(encoded) > len(text):
        result.extend(encoded[len(text):])
    return "".join(result)


def _pdf_cell(pdf, w, h, txt="", border=0, ln=0, align="", fill=False):
    """pdf.cell wrapper that sanitises text before writing."""
    pdf.cell(w, h, sanitize_pdf_text(txt), border=border, ln=ln,
             align=align, fill=fill)


def _pdf_multi_cell(pdf, w, h, txt="", border=0, align="", fill=False):
    """pdf.multi_cell wrapper that sanitises text before writing."""
    pdf.multi_cell(w, h, sanitize_pdf_text(txt), border=border,
                   align=align, fill=fill)


def get_patterns_text(df: pd.DataFrame) -> list:
    items = []
    if "platform" in df.columns and "cpa" in df.columns:
        plat_cpa = df.groupby("platform")["cpa"].mean().dropna()
        if not plat_cpa.empty:
            best = plat_cpa.idxmin()
            worst = plat_cpa.idxmax()
            items.append(f"{best} has the lowest avg CPA at {fmt_currency(plat_cpa[best])} — best platform for efficient paid starts.")
            if best != worst:
                items.append(f"{worst} has the highest avg CPA at {fmt_currency(plat_cpa[worst])} — review spend allocation.")
    if "platform" in df.columns and "paid_starts" in df.columns:
        plat_paid = df.groupby("platform")["paid_starts"].sum().dropna()
        if not plat_paid.empty:
            best = plat_paid.idxmax()
            items.append(f"{best} drives the most paid starts in total ({fmt_num(plat_paid[best])}).")
    if "format_concept" in df.columns and "cpa" in df.columns:
        concept_cpa = df.groupby("format_concept")["cpa"].mean().dropna()
        if not concept_cpa.empty:
            best_c = concept_cpa.idxmin()
            short = best_c[:60] + "…" if len(best_c) > 60 else best_c
            items.append(f'Best concept by CPA: "{short}" at {fmt_currency(concept_cpa[best_c])}.')
    if "platform" in df.columns and "cpa" in df.columns:
        plat_cpa = df.groupby("platform")["cpa"].mean().dropna()
        rows = " | ".join(f"{p}: {fmt_currency(v)}" for p, v in plat_cpa.items())
        items.append(f"Average CPA by platform — {rows}")
    if "length_s" in df.columns and "cpa" in df.columns:
        d = df.dropna(subset=["length_s", "cpa"]).copy()
        d["bucket"] = d["length_s"].apply(length_bucket)
        bucket_cpa = d.groupby("bucket")["cpa"].mean().dropna()
        ORDER = ["Under 20s", "20–30s", "31–40s", "Over 40s"]
        bucket_cpa = bucket_cpa.reindex([b for b in ORDER if b in bucket_cpa.index])
        if not bucket_cpa.empty:
            rows = " | ".join(f"{b}: {fmt_currency(v)}" for b, v in bucket_cpa.items())
            items.append(f"Average CPA by length — {rows}")
            items.append(f"{bucket_cpa.idxmin()} creatives have the lowest average CPA.")
    return items


def get_recommendations_text(df: pd.DataFrame) -> list:
    recs = []
    if "platform" in df.columns and "thumbstop_rate" in df.columns and "cpa" in df.columns:
        plat = df.groupby("platform").agg(
            thumbstop=("thumbstop_rate", "mean"),
            cpa=("cpa", "mean"),
        ).dropna()
        has_tiktok = "TikTok" in plat.index
        has_meta = "Meta" in plat.index
        if has_tiktok and has_meta:
            if plat.loc["TikTok", "thumbstop"] > plat.loc["Meta", "thumbstop"] and plat.loc["TikTok", "cpa"] > plat.loc["Meta", "cpa"]:
                recs.append("TikTok hooks are stopping scrolls but not converting. Test stronger product proof, clearer CTA, or a harder offer in the last 5 seconds.")
            if plat.loc["Meta", "thumbstop"] < plat.loc["TikTok", "thumbstop"] and plat.loc["Meta", "cpa"] < plat.loc["TikTok", "cpa"]:
                recs.append("Meta converts better despite lower thumbstop. Test more scroll-stopping hooks on Meta while keeping the same proven offer.")
    if "ctr" in df.columns and "cpa" in df.columns:
        med_ctr = df["ctr"].median()
        med_cpa = df["cpa"].median()
        hchp = df[(df["ctr"] > med_ctr) & (df["cpa"] > med_cpa)]
        if not hchp.empty:
            ids = ", ".join(hchp["creative_id"].astype(str).tolist())
            recs.append(f"Creatives {ids} have strong CTR but high CPA — drop-off is likely in the landing page or trial-to-paid flow. Test a more direct offer page.")
    if "thumbstop_rate" in df.columns and "ctr" in df.columns:
        med_t = df["thumbstop_rate"].median()
        med_c = df["ctr"].median()
        htlc = df[(df["thumbstop_rate"] > med_t) & (df["ctr"] < med_c)]
        if not htlc.empty:
            ids = ", ".join(htlc["creative_id"].astype(str).tolist())
            recs.append(f"Creatives {ids} stop the scroll but don't earn the click. Test a clearer value proposition in the first 3 seconds.")
    if "format_concept" in df.columns and "cpa" in df.columns:
        concept_cpa = df.groupby("format_concept")["cpa"].mean().dropna()
        if not concept_cpa.empty and len(concept_cpa) > 1:
            best_c = concept_cpa.idxmin()
            short = best_c[:60] + "…" if len(best_c) > 60 else best_c
            pct = (concept_cpa.mean() - concept_cpa.min()) / concept_cpa.mean() * 100
            recs.append(f'Concept "{short}" has the lowest CPA — {pct:.0f}% below average. Remix this format across other platforms.')
    if "platform" in df.columns and "cpa" in df.columns:
        plat_cpa = df.groupby("platform")["cpa"].mean().dropna()
        if len(plat_cpa) > 1:
            worst = plat_cpa.idxmax()
            best = plat_cpa.idxmin()
            recs.append(f"Consider shifting budget from {worst} (avg CPA {fmt_currency(plat_cpa[worst])}) toward {best} (avg CPA {fmt_currency(plat_cpa[best])}).")
    return recs


# ── Decision label colours for PDF table (RGB tuples: text, fill) ────────────
_PDF_DECISION_COLORS = {
    "Scale":        ((22, 101, 52),   (220, 252, 231)),
    "Keep Testing": ((146, 64, 14),   (254, 243, 199)),
    "Fix Funnel":   ((154, 52, 18),   (255, 237, 213)),
    "Cut":          ((153, 27, 27),   (254, 226, 226)),
    "Review":       ((120, 53, 15),   (254, 249, 195)),
}


class _PDF(FPDF):
    def __init__(self, goal: str, report_date: str):
        super().__init__(orientation="P", unit="mm", format="A4")
        self._goal = goal
        self._report_date = report_date
        self.set_margins(15, 15, 15)
        self.set_auto_page_break(auto=True, margin=15)

    def header(self):
        self.set_fill_color(30, 58, 138)
        self.rect(0, 0, 210, 18, "F")
        self.set_font("Helvetica", "B", 12)
        self.set_text_color(255, 255, 255)
        self.set_xy(15, 4)
        self.cell(120, 10, "Creative Performance Analyzer", ln=False)
        self.set_font("Helvetica", "", 9)
        self.set_xy(135, 4)
        self.cell(60, 5, sanitize_pdf_text(f"Goal: {self._goal}"), ln=False, align="R")
        self.set_xy(135, 9)
        self.cell(60, 5, sanitize_pdf_text(f"Generated: {self._report_date}"), ln=False, align="R")
        self.set_text_color(0, 0, 0)
        self.set_xy(15, 22)

    def footer(self):
        self.set_y(-12)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(150, 150, 150)
        self.cell(0, 5, sanitize_pdf_text(f"Page {self.page_no()} - Creative Performance Analyzer"), align="C")
        self.set_text_color(0, 0, 0)

    def section_title(self, title: str):
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(30, 58, 138)
        self.cell(0, 7, sanitize_pdf_text(title).upper(), ln=True)
        self.set_draw_color(30, 58, 138)
        self.set_line_width(0.4)
        self.line(15, self.get_y(), 195, self.get_y())
        self.ln(3)
        self.set_text_color(0, 0, 0)
        self.set_draw_color(0, 0, 0)
        self.set_line_width(0.2)


def build_pdf_report(df: pd.DataFrame, goal: str, cpa_target: float) -> bytes:
    today = datetime.date.today().strftime("%B %d, %Y")
    pdf = _PDF(goal=goal, report_date=today)
    pdf.add_page()

    # ── 1. Summary cards ────────────────────────────────────────────────────
    total_spend = df["spend"].sum() if "spend" in df.columns else float("nan")
    total_trials = df["trial_starts"].sum() if "trial_starts" in df.columns else float("nan")
    total_paid = df["paid_starts"].sum() if "paid_starts" in df.columns else float("nan")
    blended_cpa = (total_spend / total_paid) if (pd.notna(total_spend) and total_paid) else float("nan")
    best_creative = str(df.iloc[0]["creative_id"]) if (not df.empty and "creative_id" in df.columns) else "—"
    scale_n = int((df["decision_label"] == "Scale").sum()) if "decision_label" in df.columns else 0
    cut_n = int((df["decision_label"] == "Cut").sum()) if "decision_label" in df.columns else 0

    cards = [
        ("Total Spend", fmt_currency(total_spend)),
        ("Trial Starts", fmt_num(total_trials)),
        ("Paid Starts", fmt_num(total_paid)),
        ("Blended CPA", fmt_currency(blended_cpa)),
        ("Best Creative", best_creative),
        ("Scale Candidates", str(scale_n)),
        ("Cut Candidates", str(cut_n)),
    ]

    pdf.section_title("Performance Summary")
    card_w = 180 / 7
    card_h = 14
    x0 = 15
    for i, (label, value) in enumerate(cards):
        x = x0 + i * card_w
        y = pdf.get_y()
        pdf.set_fill_color(241, 245, 249)
        pdf.set_draw_color(203, 213, 225)
        pdf.set_line_width(0.3)
        pdf.rect(x, y, card_w - 1, card_h, "FD")
        pdf.set_xy(x, y + 1)
        pdf.set_font("Helvetica", "", 6.5)
        pdf.set_text_color(100, 116, 139)
        _pdf_cell(pdf, card_w - 1, 4, label, align="C")
        pdf.set_xy(x, y + 5.5)
        pdf.set_font("Helvetica", "B", 9)
        pdf.set_text_color(15, 23, 42)
        _pdf_cell(pdf, card_w - 1, 6, value, align="C")
    pdf.ln(card_h + 4)
    pdf.set_text_color(0, 0, 0)

    # ── 2. Rankings table ────────────────────────────────────────────────────
    pdf.section_title("Creative Rankings")

    tbl_cols = [
        ("Rank",       9,  "C"),
        ("Creative ID", 30, "L"),
        ("Platform",   25, "L"),
        ("Spend ($)",  22, "R"),
        ("Paid Starts", 22, "R"),
        ("CPA ($)",    22, "R"),
        ("CTR (%)",    18, "R"),
        ("Decision",   32, "C"),
    ]
    col_keys = ["_rank", "creative_id", "platform", "spend", "paid_starts", "cpa", "ctr", "decision_label"]

    header_h = 7
    row_h = 6

    ranked = df.reset_index(drop=True).copy()
    ranked.index = ranked.index + 1

    def _render_table_header():
        pdf.set_fill_color(30, 58, 138)
        pdf.set_text_color(255, 255, 255)
        pdf.set_font("Helvetica", "B", 7.5)
        for col_label, col_w, _ in tbl_cols:
            _pdf_cell(pdf, col_w, header_h, col_label, border=0, fill=True, align="C")
        pdf.ln(header_h)
        pdf.set_text_color(0, 0, 0)

    _render_table_header()

    for i, row in ranked.iterrows():
        if pdf.get_y() > 260:
            pdf.add_page()
            _render_table_header()
        decision = str(row.get("decision_label", "")) if "decision_label" in ranked.columns else ""
        txt_rgb, fill_rgb = _PDF_DECISION_COLORS.get(decision, ((50, 50, 50), (248, 250, 252)))

        fill_color = (248, 250, 252) if i % 2 == 0 else (255, 255, 255)
        for idx, (_, col_w, align) in enumerate(tbl_cols):
            key = col_keys[idx]
            if key == "_rank":
                val = str(i)
            elif key == "decision_label":
                val = decision
            elif key not in ranked.columns:
                val = "-"
            else:
                raw = row[key]
                if pd.isna(raw):
                    val = "-"
                elif key in ("spend", "cpa"):
                    val = f"${raw:,.0f}"
                elif key in ("paid_starts",):
                    val = f"{int(raw):,}"
                elif key in ("ctr",):
                    val = f"{raw:.1f}%"
                else:
                    val = str(raw)

            if key == "decision_label" and decision:
                pdf.set_fill_color(*fill_rgb)
                pdf.set_text_color(*txt_rgb)
            else:
                pdf.set_fill_color(*fill_color)
                pdf.set_text_color(30, 41, 59)

            pdf.set_font("Helvetica", "B" if key == "decision_label" else "", 7)
            safe_val = sanitize_pdf_text(val)[:22]
            _pdf_cell(pdf, col_w, row_h, safe_val, border=0, fill=True, align=align)
        pdf.ln(row_h)

    pdf.set_text_color(0, 0, 0)
    pdf.ln(4)

    # ── 3. Chart note (no kaleido/Chrome dependency) ─────────────────────────
    pdf.section_title("Charts")
    pdf.set_font("Helvetica", "I", 8.5)
    pdf.set_text_color(100, 116, 139)
    _pdf_multi_cell(pdf, 0, 5.5, "Interactive charts are displayed in the app. Download individual chart data as CSV using the buttons in the Visuals section.", border=0)
    pdf.set_text_color(0, 0, 0)
    pdf.ln(3)

    # ── 4. Pattern analysis ──────────────────────────────────────────────────
    patterns = get_patterns_text(df)
    if patterns:
        if pdf.get_y() > 240:
            pdf.add_page()
        pdf.section_title("Pattern Analysis")
        for item in patterns:
            if pdf.get_y() > 270:
                pdf.add_page()
            pdf.set_font("Helvetica", "", 8.5)
            pdf.set_text_color(30, 41, 59)
            pdf.set_fill_color(239, 246, 255)
            pdf.set_x(15)
            pdf.set_draw_color(59, 130, 246)
            pdf.set_line_width(0.8)
            text = sanitize_pdf_text(item)
            _pdf_multi_cell(pdf, 0, 5.5, f"  {text}", border="L", fill=True)
            pdf.ln(1.5)
        pdf.set_draw_color(0, 0, 0)
        pdf.set_line_width(0.2)
        pdf.ln(2)

    # ── 5. Recommendations ───────────────────────────────────────────────────
    recs = get_recommendations_text(df)
    if recs:
        if pdf.get_y() > 240:
            pdf.add_page()
        pdf.section_title("Recommendations — What to Test Next")
        for rec in recs:
            if pdf.get_y() > 270:
                pdf.add_page()
            pdf.set_font("Helvetica", "", 8.5)
            pdf.set_text_color(30, 41, 59)
            pdf.set_fill_color(240, 253, 244)
            pdf.set_x(15)
            pdf.set_draw_color(34, 197, 94)
            pdf.set_line_width(0.8)
            text = sanitize_pdf_text(rec)
            _pdf_multi_cell(pdf, 0, 5.5, f"  >>  {text}", border="L", fill=True)
            pdf.ln(1.5)

    return bytes(pdf.output())


# ══════════════════════════════════════════════════════════════════════════════
# CHART CSV DOWNLOADS
# ══════════════════════════════════════════════════════════════════════════════

def render_chart_csv_downloads(df: pd.DataFrame):
    """Offer CSV downloads for the underlying data behind each of the 4 charts."""
    st.info(
        "Chart image export is disabled in this prototype. "
        "Download the underlying chart data as CSV using the buttons below."
    )
    dl_cols = st.columns(4)

    # 1 — CPA by Creative
    with dl_cols[0]:
        if "creative_id" in df.columns and "cpa" in df.columns:
            d = df[["creative_id", "cpa"]].dropna().sort_values("cpa").rename(
                columns={"creative_id": "Creative ID", "cpa": "CPA ($)"}
            )
            st.download_button(
                label="⬇ CPA by Creative",
                data=d.to_csv(index=False).encode("utf-8"),
                file_name="chart_cpa_by_creative.csv",
                mime="text/csv",
                use_container_width=True,
                key="dl_chart_cpa",
            )

    # 2 — Paid Starts by Creative
    with dl_cols[1]:
        if "creative_id" in df.columns and "paid_starts" in df.columns:
            cols_present = [c for c in ["creative_id", "platform", "paid_starts"] if c in df.columns]
            d = df[cols_present].dropna(subset=["paid_starts"]).sort_values(
                "paid_starts", ascending=False
            ).rename(columns={"creative_id": "Creative ID", "platform": "Platform", "paid_starts": "Paid Starts"})
            st.download_button(
                label="⬇ Paid Starts",
                data=d.to_csv(index=False).encode("utf-8"),
                file_name="chart_paid_starts.csv",
                mime="text/csv",
                use_container_width=True,
                key="dl_chart_paid",
            )

    # 3 — CTR vs CPA scatter
    with dl_cols[2]:
        if "ctr" in df.columns and "cpa" in df.columns:
            cols_present = [c for c in ["creative_id", "platform", "ctr", "cpa", "spend", "decision_label"] if c in df.columns]
            d = df[cols_present].dropna(subset=["ctr", "cpa"]).rename(columns={
                "creative_id": "Creative ID", "platform": "Platform",
                "ctr": "CTR (%)", "cpa": "CPA ($)", "spend": "Spend ($)",
                "decision_label": "Decision",
            })
            st.download_button(
                label="⬇ CTR vs CPA",
                data=d.to_csv(index=False).encode("utf-8"),
                file_name="chart_ctr_vs_cpa.csv",
                mime="text/csv",
                use_container_width=True,
                key="dl_chart_scatter",
            )

    # 4 — Platform summary
    with dl_cols[3]:
        if "platform" in df.columns:
            agg = {}
            if "cpa" in df.columns:
                agg["Avg CPA ($)"] = ("cpa", "mean")
            if "paid_starts" in df.columns:
                agg["Total Paid Starts"] = ("paid_starts", "sum")
            if "ctr" in df.columns:
                agg["Avg CTR (%)"] = ("ctr", "mean")
            if "thumbstop_rate" in df.columns:
                agg["Avg Thumbstop (%)"] = ("thumbstop_rate", "mean")
            if agg:
                d = df.groupby("platform").agg(**agg).reset_index().rename(
                    columns={"platform": "Platform"}
                )
                st.download_button(
                    label="⬇ Platform Summary",
                    data=d.to_csv(index=False).encode("utf-8"),
                    file_name="chart_platform_summary.csv",
                    mime="text/csv",
                    use_container_width=True,
                    key="dl_chart_platform",
                )


# ══════════════════════════════════════════════════════════════════════════════
# INTEGRATION PLACEHOLDER FUNCTIONS
# ══════════════════════════════════════════════════════════════════════════════

def test_meta_ads_connection(access_token: str, ad_account_id: str) -> tuple[str, str]:
    """Prototype only. A real implementation would call the Meta Marketing API."""
    # TODO: Replace with requests.get("https://graph.facebook.com/v19.0/me", ...)
    if access_token and ad_account_id:
        return "Connected", "Prototype only — no live API call made."
    return "Failed", "Access token and Ad Account ID are required."


def test_tiktok_ads_connection(access_token: str, advertiser_id: str) -> tuple[str, str]:
    """Prototype only. A real implementation would call the TikTok Marketing API."""
    # TODO: Replace with requests.get("https://business-api.tiktok.com/open_api/v1.3/...")
    if access_token and advertiser_id:
        return "Connected", "Prototype only — no live API call made."
    return "Failed", "Access token and Advertiser ID are required."


def test_airtable_connection(api_key: str, base_id: str) -> tuple[str, str]:
    """Prototype only. A real implementation would call api.airtable.com."""
    # TODO: Replace with requests.get(f"https://api.airtable.com/v0/{base_id}/...", headers=...)
    if api_key and base_id:
        return "Connected", "Prototype only — no live API call made."
    return "Failed", "API key and Base ID are required."


def test_notion_connection(api_key: str) -> tuple[str, str]:
    """Prototype only. A real implementation would call api.notion.com."""
    # TODO: Replace with requests.get("https://api.notion.com/v1/users/me", headers=...)
    if api_key:
        return "Connected", "Prototype only — no live API call made."
    return "Failed", "API key is required."


def test_slack_connection(bot_token: str) -> tuple[str, str]:
    """Prototype only. A real implementation would call slack.com/api/auth.test."""
    # TODO: Replace with requests.post("https://slack.com/api/auth.test", headers=...)
    if bot_token:
        return "Connected", "Prototype only — no live API call made."
    return "Failed", "Bot token is required."


def test_google_drive_connection(api_key: str) -> tuple[str, str]:
    """Prototype only. A real implementation would use the Google Drive API."""
    # TODO: Replace with requests.get("https://www.googleapis.com/drive/v3/about", ...)
    if api_key:
        return "Connected", "Prototype only — no live API call made."
    return "Failed", "API key is required."


def test_mixpanel_connection(api_secret: str) -> tuple[str, str]:
    """Prototype only. A real implementation would call data.mixpanel.com."""
    # TODO: Replace with requests.get("https://data.mixpanel.com/api/2.0/jql", auth=...)
    if api_secret:
        return "Connected", "Prototype only — no live API call made."
    return "Failed", "API secret is required."


def test_amplitude_connection(api_key: str) -> tuple[str, str]:
    """Prototype only. A real implementation would call amplitude.com/api/2/."""
    # TODO: Replace with requests.get("https://amplitude.com/api/2/taxonomy/event", auth=...)
    if api_key:
        return "Connected", "Prototype only — no live API call made."
    return "Failed", "API key is required."


def test_hubspot_connection(private_app_token: str) -> tuple[str, str]:
    """Prototype only. A real implementation would call api.hubapi.com."""
    # TODO: Replace with requests.get("https://api.hubapi.com/crm/v3/objects/contacts", headers=...)
    if private_app_token:
        return "Connected", "Prototype only — no live API call made."
    return "Failed", "Private app token is required."


def try_load_google_sheet(url: str) -> tuple[pd.DataFrame | None, str]:
    """Attempt to load a public Google Sheets CSV export URL into a DataFrame."""
    try:
        df = pd.read_csv(url)
        return df, "ok"
    except Exception as e:
        return None, str(e)


def _status_badge(status: str) -> str:
    colors = {
        "Connected": ("#34d399", "#1a4731"),
        "Failed": ("#f87171", "#3b0a0a"),
        "Not connected": ("#8A9BC8", "#1a2035"),
        "Prototype only": ("#fbbf24", "#3b3515"),
    }
    fg, bg = colors.get(status, ("#8A9BC8", "#1a2035"))
    return (
        f"<span style='background:{bg};color:{fg};border:1px solid {fg};"
        f"padding:2px 10px;border-radius:20px;font-size:0.78rem;font-weight:700;'>{status}</span>"
    )


def render_integration_card(
    title: str,
    description: str,
    fields: list[dict],
    possible_pulls: list[str],
    test_fn,
    state_key: str,
):
    """Render a single integration card with inputs, test button, and data pull list."""
    with st.container():
        st.markdown(
            f"<div style='background:#161B27;border:1px solid #2A3350;border-radius:10px;"
            f"padding:20px 24px;margin-bottom:16px;'>",
            unsafe_allow_html=True,
        )
        header_col, badge_col = st.columns([4, 1])
        with header_col:
            st.markdown(f"**{title}**")
            st.caption(description)
        with badge_col:
            status = st.session_state.get(f"status_{state_key}", "Not connected")
            st.markdown(_status_badge(status), unsafe_allow_html=True)

        collected = {}
        for field in fields:
            key = f"{state_key}_{field['key']}"
            if field.get("password"):
                collected[field["key"]] = st.text_input(
                    field["label"], key=key, type="password",
                    placeholder=field.get("placeholder", ""),
                )
            else:
                collected[field["key"]] = st.text_input(
                    field["label"], key=key,
                    placeholder=field.get("placeholder", ""),
                )

        if st.button(f"Test connection — {title}", key=f"btn_{state_key}"):
            result_status, result_msg = test_fn(**collected)
            st.session_state[f"status_{state_key}"] = result_status
            if result_status == "Connected":
                st.success(result_msg)
            else:
                st.error(result_msg)

        with st.expander("Possible data pulls"):
            for item in possible_pulls:
                st.markdown(f"- {item}")

        st.markdown("</div>", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE: CUSTOM COLUMN MAPPER
# ══════════════════════════════════════════════════════════════════════════════

_EXPECTED_COLS = {
    "creative_id": "Creative ID",
    "platform": "Platform",
    "spend": "Spend ($)",
    "paid_starts": "Paid Starts",
    "trial_starts": "Trial Starts",
    "ctr": "CTR (%)",
    "thumbstop_rate": "Thumbstop Rate (%)",
    "hold_6s": "6s Hold Rate (%)",
}


def render_column_mapper(df: pd.DataFrame) -> pd.DataFrame | None:
    """Show UI to map unrecognised column names. Returns remapped df or None."""
    missing = [k for k in _EXPECTED_COLS if k not in df.columns]
    if not missing:
        return None
    with st.expander(f"⚙️ Column Mapper — {len(missing)} column(s) not auto-detected", expanded=True):
        st.info(
            "Some expected columns were not found automatically. "
            "Use the dropdowns to map your column names, or leave as '(skip)' to continue without them."
        )
        options = ["(skip)"] + list(df.columns)
        mapping: dict[str, str] = {}
        num_cols = min(len(missing), 4)
        mapper_cols = st.columns(num_cols)
        for idx, col_key in enumerate(missing):
            with mapper_cols[idx % num_cols]:
                chosen = st.selectbox(
                    _EXPECTED_COLS[col_key],
                    options,
                    key=f"colmap_{col_key}",
                )
                if chosen != "(skip)":
                    mapping[chosen] = col_key
        if mapping:
            return df.rename(columns=mapping)
    return None


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE: DATE COLUMN DETECTION & CREATIVE FATIGUE
# ══════════════════════════════════════════════════════════════════════════════

_DATE_COL_NAMES = {"date", "week", "period", "month", "week_num", "day", "period_label", "run_date"}


def detect_date_col(df: pd.DataFrame) -> str | None:
    for c in df.columns:
        if c.lower() in _DATE_COL_NAMES:
            return c
    return None


def detect_fatigue_ids(df: pd.DataFrame) -> set:
    """Return creative_ids with monotonically rising CPA or falling CTR across periods."""
    date_col = detect_date_col(df)
    if date_col is None or "creative_id" not in df.columns:
        return set()
    fatigue = set()
    for cid, sub in df.groupby("creative_id"):
        sub = sub.sort_values(date_col)
        if "cpa" in sub.columns:
            vals = sub["cpa"].dropna().tolist()
            if len(vals) >= 2 and all(vals[i] < vals[i + 1] for i in range(len(vals) - 1)):
                fatigue.add(cid)
        if "ctr" in sub.columns:
            vals = sub["ctr"].dropna().tolist()
            if len(vals) >= 2 and all(vals[i] > vals[i + 1] for i in range(len(vals) - 1)):
                fatigue.add(cid)
    return fatigue


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE: STATISTICAL SIGNIFICANCE
# ══════════════════════════════════════════════════════════════════════════════

def compute_significance(df: pd.DataFrame) -> pd.DataFrame:
    """Add stat_sig column — two-proportion z-test vs. rest of cohort."""
    df = df.copy()
    if "trial_starts" not in df.columns or "paid_starts" not in df.columns:
        return df
    total_trials = df["trial_starts"].sum(skipna=True)
    total_paid = df["paid_starts"].sum(skipna=True)
    labels = []
    for _, row in df.iterrows():
        n1 = row.get("trial_starts", 0)
        x1 = row.get("paid_starts", 0)
        if pd.isna(n1) or pd.isna(x1) or n1 <= 0:
            labels.append("—"); continue
        n2 = total_trials - n1
        x2 = total_paid - x1
        if n2 <= 0 or (x1 + x2) <= 0:
            labels.append("—"); continue
        p_pool = (x1 + x2) / (n1 + n2)
        denom = math.sqrt(p_pool * (1 - p_pool) * (1 / n1 + 1 / n2)) if p_pool not in (0, 1) else 0
        if denom == 0:
            labels.append("—"); continue
        z = abs((x1 / n1) - (x2 / n2)) / denom
        if z > 2.576:   labels.append("✓ 99% sig.")
        elif z > 1.96:  labels.append("✓ 95% sig.")
        elif z > 1.645: labels.append("~ 90% sig.")
        else:           labels.append("Not sig.")
    df["stat_sig"] = labels
    return df


def render_significance(df: pd.DataFrame):
    if "stat_sig" not in df.columns:
        st.info("Statistical significance requires Trial Starts + Paid Starts columns.")
        return
    st.caption(
        "Two-proportion z-test: compares each creative's Trial→Paid CVR against the rest of the cohort. "
        "✓ 95% means we're 95% confident this creative's conversion rate is genuinely different."
    )
    show_cols = [c for c in ["creative_id", "platform", "trial_starts", "paid_starts", "cpa", "stat_sig"] if c in df.columns]
    d = df[show_cols].copy().reset_index(drop=True)
    d.index = d.index + 1
    d.index.name = "Rank"
    d = d.rename(columns={
        "creative_id": "Creative", "platform": "Platform",
        "trial_starts": "Trial Starts", "paid_starts": "Paid Starts",
        "cpa": "CPA ($)", "stat_sig": "Significance",
    })

    def _color_sig(val):
        v = str(val)
        if "99%" in v: return "color:#34d399;font-weight:bold"
        if "95%" in v: return "color:#4F8EF7;font-weight:bold"
        if "90%" in v: return "color:#fbbf24"
        if "Not" in v:  return "color:#8A9BC8"
        return ""

    styler = d.style.map(_color_sig, subset=["Significance"]) if "Significance" in d.columns else d.style
    st.dataframe(styler, use_container_width=True)


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE: BUDGET REALLOCATION
# ══════════════════════════════════════════════════════════════════════════════

def render_budget_reallocation(df: pd.DataFrame):
    if not {"spend", "cpa", "creative_id"}.issubset(df.columns):
        st.info("Need Spend + CPA + Creative ID to generate a reallocation plan.")
        return
    valid = df[["creative_id", "spend", "cpa"] +
               ([c] for c in ["paid_starts", "decision_label"] if c in df.columns and False)].copy()
    valid = df[[c for c in ["creative_id", "spend", "cpa", "paid_starts", "decision_label"] if c in df.columns]].dropna(subset=["cpa", "spend"]).copy()
    valid = valid[valid["cpa"] > 0]
    if len(valid) < 2:
        st.info("Need at least 2 creatives with spend + CPA data.")
        return
    total_budget = valid["spend"].sum()
    valid["efficiency"] = 1.0 / valid["cpa"]
    total_eff = valid["efficiency"].sum()
    valid["rec_spend"] = (valid["efficiency"] / total_eff * total_budget).round(0)
    valid["delta"] = (valid["rec_spend"] - valid["spend"]).round(0)
    valid["exp_paid"] = (valid["rec_spend"] / valid["cpa"]).round(0)
    st.caption(
        f"Same total budget (${total_budget:,.0f}) reallocated proportionally to 1/CPA — "
        "lower CPA gets more. Green = increase, red = decrease."
    )
    show_cols = ["creative_id", "spend", "rec_spend", "delta", "cpa", "exp_paid"] + \
                (["decision_label"] if "decision_label" in valid.columns else [])
    display = valid[show_cols].copy().reset_index(drop=True)
    display.index = display.index + 1
    display.columns = ["Creative", "Current Spend ($)", "Rec. Spend ($)", "Delta ($)", "CPA ($)", "Exp. Paid Starts"] + \
                      (["Decision"] if "decision_label" in valid.columns else [])

    def _color_delta(v):
        try:
            return "color:#34d399" if float(v) > 0 else ("color:#f87171" if float(v) < 0 else "")
        except Exception:
            return ""

    styler = display.style.map(_color_delta, subset=["Delta ($)"])
    st.dataframe(styler, use_container_width=True)
    st.download_button(
        "⬇ Download Reallocation Plan as CSV",
        data=display.to_csv(index=False).encode("utf-8"),
        file_name="budget_reallocation.csv", mime="text/csv", key="dl_realloc",
    )


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE: SPEND PACING
# ══════════════════════════════════════════════════════════════════════════════

def render_spend_pacing(df: pd.DataFrame):
    if "spend" not in df.columns:
        st.info("Need a Spend column for pacing.")
        return
    p1, p2 = st.columns(2)
    with p1:
        monthly_budget = st.number_input(
            "Monthly Budget ($)", min_value=0.0, value=10000.0, step=500.0, key="pacing_budget"
        )
    with p2:
        days_elapsed = st.number_input(
            "Days elapsed in month", min_value=1, max_value=31,
            value=min(datetime.date.today().day, 31), key="pacing_days",
        )
    days_in_month = 30
    paced = monthly_budget * days_elapsed / days_in_month
    total_spend = df["spend"].sum(skipna=True)
    pacing_pct = (total_spend / paced * 100) if paced > 0 else 0
    delta_v = total_spend - paced
    status = "over-pacing" if delta_v > 0 else "under-pacing"
    m1, m2, m3, m4 = st.columns(4)
    m1.metric("Total Spend to Date", f"${total_spend:,.0f}")
    m2.metric("Paced Budget", f"${paced:,.0f}")
    m3.metric("vs. Pace", f"${abs(delta_v):,.0f} {status}")
    m4.metric("Pacing Rate", f"{pacing_pct:.1f}%")
    if "creative_id" in df.columns:
        pc = df.groupby("creative_id")["spend"].sum(numeric_only=True).reset_index()
        pc.columns = ["Creative", "Spend to Date ($)"]
        pc["Share (%)"] = (pc["Spend to Date ($)"] / total_spend * 100).round(1)
        pc["Projected Month-End ($)"] = (pc["Spend to Date ($)"] / days_elapsed * days_in_month).round(0)
        st.dataframe(
            pc.sort_values("Spend to Date ($)", ascending=False).reset_index(drop=True),
            use_container_width=True,
        )


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE: BENCHMARK COMPARISON
# ══════════════════════════════════════════════════════════════════════════════

def render_benchmark_comparison(df: pd.DataFrame, benchmarks: dict):
    if "platform" not in df.columns:
        st.info("Need a Platform column for benchmark comparison.")
        return
    rows = []
    for plat, grp in df.groupby("platform"):
        row: dict = {"Platform": plat}
        if "ctr" in grp.columns:
            avg = grp["ctr"].mean(skipna=True)
            bench = benchmarks.get(f"{plat}_ctr")
            row["Your Avg CTR (%)"] = round(avg, 2) if pd.notna(avg) else "—"
            if bench:
                row["Benchmark CTR (%)"] = bench
                if pd.notna(avg):
                    d = avg - bench
                    row["CTR Delta"] = f"+{d:.2f}%" if d >= 0 else f"{d:.2f}%"
        if "cpa" in grp.columns:
            avg = grp["cpa"].mean(skipna=True)
            bench = benchmarks.get(f"{plat}_cpa")
            row["Your Avg CPA ($)"] = round(avg, 2) if pd.notna(avg) else "—"
            if bench:
                row["Benchmark CPA ($)"] = bench
                if pd.notna(avg):
                    d = avg - bench
                    row["CPA Delta"] = f"+${d:.0f} above" if d > 0 else f"${-d:.0f} below"
        if "thumbstop_rate" in grp.columns:
            avg = grp["thumbstop_rate"].mean(skipna=True)
            bench = benchmarks.get(f"{plat}_thumbstop")
            row["Your Avg Thumbstop (%)"] = round(avg, 2) if pd.notna(avg) else "—"
            if bench:
                row["Benchmark Thumbstop (%)"] = bench
        rows.append(row)
    if rows:
        st.caption("Set platform benchmark targets in Advanced Settings → Benchmark Targets to unlock delta columns.")
        st.dataframe(pd.DataFrame(rows), use_container_width=True)
    else:
        st.info("No platform data available.")


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE: TREND SPARKLINES
# ══════════════════════════════════════════════════════════════════════════════

def render_trend_sparklines(df: pd.DataFrame):
    date_col = detect_date_col(df)
    if date_col is None:
        st.info(
            "No date/period column detected. Add a column named 'Week', 'Date', or 'Period' "
            "to your data to see per-creative trend sparklines."
        )
        return
    if "creative_id" not in df.columns:
        st.info("Need a Creative ID column for sparklines.")
        return
    metric_options = [m for m in ["cpa", "ctr", "paid_starts", "thumbstop_rate"] if m in df.columns]
    if not metric_options:
        st.info("No metric columns available for sparklines.")
        return
    metric_labels = {
        "cpa": "CPA ($)", "ctr": "CTR (%)",
        "paid_starts": "Paid Starts", "thumbstop_rate": "Thumbstop (%)",
    }
    sel_metric = st.selectbox(
        "Metric to trend",
        metric_options,
        format_func=lambda m: metric_labels[m],
        key="sparkline_metric",
    )
    creatives = [
        c for c in df["creative_id"].unique()
        if df[df["creative_id"] == c][sel_metric].dropna().shape[0] >= 2
    ]
    if not creatives:
        st.info("Need at least 2 data points per creative to draw a trend. Check your date/period column.")
        return
    cols = st.columns(min(len(creatives), 4))
    for i, cid in enumerate(creatives):
        sub = df[df["creative_id"] == cid].sort_values(date_col)
        with cols[i % 4]:
            fig = px.line(
                sub, x=date_col, y=sel_metric, markers=True,
                labels={sel_metric: "", date_col: ""},
                title=str(cid),
            )
            fig.update_layout(
                height=160, margin=dict(t=30, b=5, l=10, r=10),
                paper_bgcolor="#161B27", plot_bgcolor="#161B27",
                font_color="#FAFAFA", showlegend=False,
                xaxis=dict(showticklabels=False, gridcolor="#1E2A45"),
                yaxis=dict(gridcolor="#1E2A45"),
            )
            fig.update_traces(line_color="#4F8EF7", marker_color="#34d399")
            st.plotly_chart(fig, use_container_width=True)


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE: SUMMARY SNIPPET
# ══════════════════════════════════════════════════════════════════════════════

def build_summary_snippet(df: pd.DataFrame, goal: str, cpa_target: float) -> str:
    lines = [
        f"Creative Performance Summary — {datetime.date.today().strftime('%B %d, %Y')}",
        f"Goal: {goal}  |  CPA Target: ${cpa_target:,.0f}",
        "",
    ]
    if not df.empty and "creative_id" in df.columns:
        top = df.iloc[0]
        lines.append(f"Top creative: {top['creative_id']}")
        if "cpa" in df.columns and pd.notna(top.get("cpa")):
            lines.append(f"  - CPA: ${top['cpa']:,.2f}")
        if "paid_starts" in df.columns and pd.notna(top.get("paid_starts")):
            lines.append(f"  - Paid Starts: {int(top['paid_starts']):,}")
        if "decision_label" in df.columns:
            lines.append(f"  - Action: {top.get('decision_label', '')}")
        lines.append("")
    if "decision_label" in df.columns:
        for label in ["Scale", "Keep Testing", "Fix Funnel", "Cut"]:
            ids = df[df["decision_label"] == label]["creative_id"].tolist()
            if ids:
                lines.append(f"{label}: {', '.join(str(c) for c in ids)}")
        lines.append("")
    if "spend" in df.columns:
        lines.append(f"Total Spend: ${df['spend'].sum(skipna=True):,.0f}")
    if "paid_starts" in df.columns:
        total_paid = df["paid_starts"].sum(skipna=True)
        lines.append(f"Total Paid Starts: {int(total_paid):,}")
        if "spend" in df.columns and total_paid > 0:
            lines.append(f"Blended CPA: ${df['spend'].sum(skipna=True) / total_paid:,.2f}")
    lines += ["", "Generated by Creative Performance Analyzer"]
    return "\n".join(lines)


# ══════════════════════════════════════════════════════════════════════════════
# FEATURE: POWERPOINT EXPORT
# ══════════════════════════════════════════════════════════════════════════════

def build_pptx_report(df: pd.DataFrame, goal: str, cpa_target: float) -> bytes:
    from pptx import Presentation as _Prs
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = _Prs()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def _txt(slide, left, top, width, height, text, size=13, bold=False,
             color=(250, 250, 250), align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = sanitize_pdf_text(str(text))
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        p.alignment = align

    def _bg(slide, color=(0x0E, 0x11, 0x17)):
        bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(*color)
        bg.line.fill.background()

    # ── Slide 1: Cover ────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    _bg(s1)
    _txt(s1, 1.5, 2.0, 10, 1.4, "Creative Performance Report",
         size=36, bold=True, color=(0x4F, 0x8E, 0xF7), align=PP_ALIGN.CENTER)
    _txt(s1, 1.5, 3.6, 10, 0.7,
         f"Goal: {goal}  |  CPA Target: ${cpa_target:,.0f}  |  {datetime.date.today().strftime('%B %d, %Y')}",
         size=14, color=(0x8A, 0x9B, 0xC8), align=PP_ALIGN.CENTER)

    # ── Slide 2: Summary metrics ──────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    _bg(s2)
    _txt(s2, 0.5, 0.2, 12, 0.7, "Performance Summary", size=22, bold=True, color=(0x4F, 0x8E, 0xF7))
    metrics: list[tuple[str, str]] = []
    if "spend" in df.columns:
        metrics.append(("Total Spend", f"${df['spend'].sum(skipna=True):,.0f}"))
    if "paid_starts" in df.columns:
        tp = df["paid_starts"].sum(skipna=True)
        metrics.append(("Paid Starts", f"{int(tp):,}"))
        if "spend" in df.columns and tp > 0:
            metrics.append(("Blended CPA", f"${df['spend'].sum(skipna=True)/tp:,.2f}"))
    if not df.empty and "creative_id" in df.columns:
        metrics.append(("Top Creative", str(df.iloc[0]["creative_id"])))
    if "decision_label" in df.columns:
        metrics.append(("Scale", str((df["decision_label"] == "Scale").sum())))
        metrics.append(("Cut", str((df["decision_label"] == "Cut").sum())))
    cw = 1.8
    for i, (lbl, val) in enumerate(metrics[:6]):
        x = 0.5 + i * (cw + 0.22)
        box = s2.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(cw), Inches(1.5))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x27)
        box.line.color.rgb = RGBColor(0x2A, 0x33, 0x50)
        _txt(s2, x + 0.05, 1.28, cw - 0.1, 0.4, lbl, size=8, color=(0x8A, 0x9B, 0xC8), align=PP_ALIGN.CENTER)
        _txt(s2, x + 0.05, 1.65, cw - 0.1, 0.85, val, size=17, bold=True, color=(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER)

    # ── Slide 3: Rankings table ───────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    _bg(s3)
    _txt(s3, 0.5, 0.2, 12, 0.7, "Creative Rankings", size=22, bold=True, color=(0x4F, 0x8E, 0xF7))
    tbl_keys = [c for c in ["creative_id","platform","spend","paid_starts","cpa","ctr","decision_label"] if c in df.columns]
    tbl_hdrs = {"creative_id":"Creative","platform":"Platform","spend":"Spend ($)","paid_starts":"Paid Starts",
                "cpa":"CPA ($)","ctr":"CTR (%)","decision_label":"Decision"}
    top10 = df.head(10).reset_index(drop=True)
    if tbl_keys and not top10.empty:
        tbl = s3.shapes.add_table(
            len(top10) + 1, len(tbl_keys), Inches(0.4), Inches(1.0), Inches(12.5), Inches(6.0)
        ).table
        for j, k in enumerate(tbl_keys):
            c = tbl.cell(0, j)
            c.text = tbl_hdrs.get(k, k)
            c.text_frame.paragraphs[0].font.bold = True
            c.text_frame.paragraphs[0].font.size = Pt(9)
            c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            c.fill.solid(); c.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x8A)
        for i, (_, row) in enumerate(top10.iterrows()):
            for j, k in enumerate(tbl_keys):
                c = tbl.cell(i + 1, j)
                v = row.get(k, "")
                if k in ("spend", "cpa") and pd.notna(v):
                    c.text = f"${float(v):,.0f}"
                elif k == "ctr" and pd.notna(v):
                    c.text = f"{float(v):.1f}%"
                elif k == "paid_starts" and pd.notna(v):
                    c.text = f"{int(v):,}"
                else:
                    c.text = sanitize_pdf_text(str(v)) if pd.notna(v) else "-"
                c.text_frame.paragraphs[0].font.size = Pt(8)
                c.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1E, 0x29, 0x3B)
                bg = (0xF1, 0xF5, 0xF9) if i % 2 == 0 else (0xFF, 0xFF, 0xFF)
                c.fill.solid(); c.fill.fore_color.rgb = RGBColor(*bg)

    # ── Slide 4: Recommendations ──────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    _bg(s4)
    _txt(s4, 0.5, 0.2, 12, 0.7, "Recommendations", size=22, bold=True, color=(0x4F, 0x8E, 0xF7))
    recs = get_recommendations_text(df)
    box = s4.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(12.3), Inches(6.0))
    tf4 = box.text_frame
    tf4.word_wrap = True
    for i, rec in enumerate(recs[:14]):
        p = tf4.paragraphs[0] if i == 0 else tf4.add_paragraph()
        p.text = f"- {sanitize_pdf_text(rec)}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
        p.space_after = Pt(4)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ══════════════════════════════════════════════════════════════════════════════
# PAGE RENDERERS
# ══════════════════════════════════════════════════════════════════════════════

def page_analyzer(
    cpa_target: float = CPA_TARGET,
    ce_weights: dict | None = None,
    ffq_weights: dict | None = None,
):
    if ce_weights is None:
        ce_weights = CE_WEIGHTS_DEFAULT
    if ffq_weights is None:
        ffq_weights = FFQ_WEIGHTS_DEFAULT

    st.markdown(
        """
        <div style="display:flex;align-items:center;gap:12px;margin-bottom:0.2rem;">
          <span style="font-size:2rem;">📊</span>
          <span style="font-size:1.7rem;font-weight:800;letter-spacing:-0.02em;">Creative Performance Analyzer</span>
        </div>
        <p style="color:#8A9BC8;margin-top:0;margin-bottom:1rem;font-size:0.95rem;">
          Upload your ad creative data · Choose a campaign goal · Review rankings · Read recommendations
        </p>
        """,
        unsafe_allow_html=True,
    )
    st.divider()

    upload_col, sample_col = st.columns([3, 1])
    with upload_col:
        uploaded_file = st.file_uploader(
            "Upload CSV or Excel file",
            type=["csv", "xlsx", "xls"],
            help="Columns like Creative ID, Platform, Spend, Thumbstop Rate, CTR, Trial Starts, Paid Starts",
            key=f"uploader_{st.session_state.uploader_key}",
        )
    with sample_col:
        st.write("")
        st.write("")
        if st.button("📂  Use sample LoopNote data", use_container_width=True):
            df_raw, warns = load_and_clean(SAMPLE_CSV)
            st.session_state.df_raw = df_raw
            st.session_state.warnings = warns

    # Also accept data loaded from Google Sheets on Integrations page
    if uploaded_file is not None:
        try:
            df_raw, warns = load_and_clean(uploaded_file)
            st.session_state.df_raw = df_raw
            st.session_state.warnings = warns
        except Exception as e:
            st.error(f"Could not read file: {e}")
    elif st.session_state.get("sheets_df") is not None and st.session_state.df_raw is None:
        st.info("Using data loaded from Google Sheets (Integrations page).")
        st.session_state.df_raw = st.session_state.sheets_df

    if st.session_state.df_raw is not None:
        df_raw = st.session_state.df_raw

        # ── Column mapper: let users fix unmapped columns ─────────────────────
        mapped = render_column_mapper(df_raw)
        if mapped is not None:
            st.session_state.df_raw = mapped
            df_raw = mapped

        for w in st.session_state.get("warnings", []):
            st.warning(w)

        df, kpi_warnings = calculate_kpis(df_raw)
        df = assign_decision_labels(df, cpa_target=cpa_target)

        # Flag fatiguing creatives if a date column is present
        fatigue_ids = detect_fatigue_ids(df)
        if fatigue_ids and "decision_label" in df.columns:
            df.loc[df["creative_id"].isin(fatigue_ids), "decision_label"] = "Fatiguing"

        goal_options = [
            "Paid Starts", "Trial Starts", "Efficient Paid Starts",
            "Efficient Trial Starts", "Creative Engagement", "Full Funnel Quality",
        ]
        _default_goal_idx = 0
        if st.session_state.get("loaded_goal") in goal_options:
            _default_goal_idx = goal_options.index(st.session_state.loaded_goal)
            st.session_state.loaded_goal = None

        goal_col, spacer, clear_col = st.columns([2, 4, 1])
        with goal_col:
            goal = st.selectbox("Campaign Goal", goal_options, index=_default_goal_idx)
        with clear_col:
            st.write("")  # vertical alignment nudge
            if st.button("🗑 Clear analysis", use_container_width=True, key="btn_clear_analysis"):
                for _k in ("df_raw", "sheets_df", "warnings", "loaded_goal"):
                    st.session_state[_k] = None
                st.session_state.uploader_key += 1
                st.rerun()

        df_ranked = rank_by_goal(df, goal, cpa_target=cpa_target, ce_weights=ce_weights, ffq_weights=ffq_weights)
        df_ranked = assign_decision_labels(df_ranked, cpa_target=cpa_target)

        # ── Save this analysis ────────────────────────────────────────────────
        with st.expander("💾 Save this analysis", expanded=False):
            save_name_col, save_btn_col = st.columns([4, 1])
            with save_name_col:
                save_name = st.text_input(
                    "Analysis name",
                    placeholder="e.g. Week 12 LoopNote Review",
                    label_visibility="collapsed",
                    key="save_name_input",
                )
            with save_btn_col:
                if st.button("Save", use_container_width=True, key="btn_save_analysis"):
                    if save_name.strip():
                        save_analysis(save_name, df_raw, goal, cpa_target)
                        st.success(f"Saved as \"{save_name.strip()}\"")
                        st.rerun()
                    else:
                        st.warning("Enter a name before saving.")

        st.divider()

        st.markdown("<div class='section-header'>Performance Summary</div>", unsafe_allow_html=True)
        render_summary_cards(df_ranked, goal)

        if kpi_warnings:
            with st.expander("⚠️  Some KPIs could not be calculated — expand for details"):
                for w in kpi_warnings:
                    st.caption(f"• {w}")

        st.divider()

        if goal == "Efficient Paid Starts" and "scale_candidate" in df_ranked.columns:
            candidates = df_ranked[df_ranked["scale_candidate"] == True]["creative_id"].tolist()
            if candidates:
                st.success(f"🚀  **Scale Candidates** (CPA < ${cpa_target:,.0f}): {', '.join(str(c) for c in candidates)}")

        rank_hdr_col, csv_btn_col = st.columns([5, 1])
        with rank_hdr_col:
            st.markdown(f"<div class='section-header'>Creative Rankings — {goal}</div>", unsafe_allow_html=True)
        with csv_btn_col:
            csv_bytes = build_export_csv(df_ranked)
            st.download_button(
                label="⬇ Download Rankings as CSV",
                data=csv_bytes,
                file_name=f"creative_rankings_{goal.lower().replace(' ', '_')}.csv",
                mime="text/csv",
                use_container_width=True,
                key="dl_csv",
            )
        render_ranking_table(df_ranked)

        # ── Statistical significance ───────────────────────────────────────────
        df_ranked_sig = compute_significance(df_ranked)
        with st.expander("🧪 Statistical Significance", expanded=False):
            render_significance(df_ranked_sig)

        # ── Fatigue alert ─────────────────────────────────────────────────────
        ranked_fatigue = detect_fatigue_ids(df_ranked)
        if ranked_fatigue:
            st.warning(
                f"⚠️ **Creative Fatigue detected** in: {', '.join(str(c) for c in ranked_fatigue)}. "
                "These creatives show consistently rising CPA or falling CTR across periods. "
                "Consider refreshing creative or testing new angles."
            )
        elif detect_date_col(df_ranked) is not None:
            st.success("✅ No monotonic fatigue patterns detected across periods.")

        st.divider()

        st.markdown("<div class='section-header'>Visuals</div>", unsafe_allow_html=True)
        render_charts(df_ranked, cpa_target=cpa_target)
        render_chart_csv_downloads(df_ranked)

        # ── Benchmark comparison ──────────────────────────────────────────────
        with st.expander("📊 Platform Benchmark Comparison", expanded=False):
            render_benchmark_comparison(df_ranked, st.session_state.get("_benchmarks", {}))

        # ── Trend sparklines ──────────────────────────────────────────────────
        with st.expander("📈 Creative Trend Sparklines", expanded=False):
            render_trend_sparklines(df_ranked)

        st.divider()

        st.markdown("<div class='section-header'>What patterns are emerging?</div>", unsafe_allow_html=True)
        render_patterns(df_ranked)

        st.divider()

        st.markdown("<div class='section-header'>What should we test next?</div>", unsafe_allow_html=True)
        render_recommendations(df_ranked)

        # ── Budget reallocation ───────────────────────────────────────────────
        with st.expander("💰 Budget Reallocation Plan", expanded=False):
            render_budget_reallocation(df_ranked)

        # ── Spend pacing ──────────────────────────────────────────────────────
        with st.expander("📅 Spend Pacing Tracker", expanded=False):
            render_spend_pacing(df_ranked)

        st.divider()

        # ── Summary snippet ───────────────────────────────────────────────────
        with st.expander("📋 Summary Snippet (copy for Slack / email)", expanded=False):
            snippet = build_summary_snippet(df_ranked, goal, cpa_target)
            st.text_area(
                "Copy this summary",
                value=snippet,
                height=220,
                label_visibility="collapsed",
                key="summary_snippet_area",
            )

        # ── Exports ───────────────────────────────────────────────────────────
        st.markdown("<div class='section-header'>Export Reports</div>", unsafe_allow_html=True)
        export_col1, export_col2 = st.columns(2)
        today_str = datetime.date.today().strftime("%Y-%m-%d")
        safe_goal = re.sub(r"[^\w]+", "_", goal.lower()).strip("_")

        with export_col1:
            try:
                pdf_bytes = build_pdf_report(df_ranked, goal, cpa_target)
                st.download_button(
                    label="⬇ Download PDF Report",
                    data=pdf_bytes,
                    file_name=f"creative_report_{safe_goal}_{today_str}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    key="dl_pdf",
                )
            except Exception as _pdf_err:
                st.error(f"PDF export failed: {_pdf_err}")

        with export_col2:
            try:
                pptx_bytes = build_pptx_report(df_ranked, goal, cpa_target)
                st.download_button(
                    label="⬇ Download PowerPoint Report",
                    data=pptx_bytes,
                    file_name=f"creative_report_{safe_goal}_{today_str}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    key="dl_pptx",
                )
            except Exception as _pptx_err:
                st.error(f"PPTX export failed: {_pptx_err}")

    else:
        st.markdown(
            """
            <div style="background:#161B27;border:1px solid #2A3350;border-radius:12px;
                        padding:32px 40px;max-width:700px;margin:24px auto 0;">
              <h3 style="margin-top:0;color:#FAFAFA;">How it works</h3>
              <div style="display:grid;grid-template-columns:40px 1fr;gap:8px 12px;align-items:start;font-size:0.95rem;color:#C5CFDF;">
                <span style="color:#4F8EF7;font-weight:700;font-size:1.1rem;">1</span>
                <span><b>Upload data</b> — CSV or Excel with your ad creative metrics, or connect Google Sheets via the Integrations page, or use the built-in LoopNote sample.</span>
                <span style="color:#4F8EF7;font-weight:700;font-size:1.1rem;">2</span>
                <span><b>Choose a campaign goal</b> — Paid Starts, Trial Starts, Efficient CPA, Creative Engagement, or Full Funnel Quality.</span>
                <span style="color:#4F8EF7;font-weight:700;font-size:1.1rem;">3</span>
                <span><b>Review rankings</b> — Creatives are scored and labelled: Scale, Keep Testing, Fix Funnel, Cut, or Review.</span>
                <span style="color:#4F8EF7;font-weight:700;font-size:1.1rem;">4</span>
                <span><b>Read recommendations</b> — Rule-based insights tell you exactly what to test next.</span>
                <span style="color:#4F8EF7;font-weight:700;font-size:1.1rem;">5</span>
                <span><b>Explore integrations</b> — See how this tool connects to Meta Ads, TikTok, Google Sheets, Airtable, Slack, and more via the Integrations page.</span>
                <span style="color:#4F8EF7;font-weight:700;font-size:1.1rem;">6</span>
                <span><b>Understand the methodology</b> — Full KPI formulas, ranking logic, and decision-label criteria are documented on the Methodology page.</span>
              </div>
              <p style="margin-bottom:0;margin-top:20px;color:#8A9BC8;font-size:0.88rem;">
                Supported columns: Creative ID · Platform · Format/Concept · Length · Spend · Thumbstop Rate ·
                6s Hold Rate · CTR · Trial Starts · Paid Starts
              </p>
            </div>
            """,
            unsafe_allow_html=True,
        )


def page_integrations():
    st.markdown(
        """
        <div style="display:flex;align-items:center;gap:12px;margin-bottom:0.2rem;">
          <span style="font-size:2rem;">🔌</span>
          <span style="font-size:1.7rem;font-weight:800;letter-spacing:-0.02em;">Integrations</span>
        </div>
        <p style="color:#8A9BC8;margin-top:0;margin-bottom:0.5rem;font-size:0.95rem;">
          Connect your data sources to turn the analyzer into a lightweight creative intelligence hub.
        </p>
        """,
        unsafe_allow_html=True,
    )

    st.warning(
        "**Security notice:** This is a prototype integration layer. API keys entered here are stored only in your "
        "browser session and are never saved or logged. Production versions should use OAuth, encrypted secret storage, "
        "environment variables, and role-based permissions."
    )

    st.divider()

    # ── Why these integrations matter ────────────────────────────────────────
    st.markdown("<div class='section-header'>Why these integrations matter</div>", unsafe_allow_html=True)
    st.markdown(
        "<div class='insight-box'>"
        "Manual CSV upload is useful for quick analysis, but growth teams usually make creative decisions from scattered data. "
        "Ad platforms show performance, Sheets and Airtable track test plans, Drive stores assets, Slack and Notion capture "
        "team learnings, and product analytics tools show whether paid users actually activate. This integration layer shows "
        "how the analyzer could become a lightweight creative intelligence hub instead of a one-off spreadsheet review."
        "</div>",
        unsafe_allow_html=True,
    )

    st.divider()

    # ── Google Sheets (real data pull) ───────────────────────────────────────
    st.markdown("<div class='section-header'>Google Sheets — Live data pull</div>", unsafe_allow_html=True)
    st.caption("Paste a public Google Sheets CSV export URL to load data directly into the Analyzer.")

    sheets_url = st.text_input(
        "Google Sheets CSV export URL",
        key="gs_url",
        placeholder="https://docs.google.com/spreadsheets/d/.../export?format=csv",
    )

    gs_status = st.session_state.get("status_google_sheets", "Not connected")
    st.markdown(_status_badge(gs_status), unsafe_allow_html=True)

    if st.button("Load from Google Sheets", key="btn_google_sheets"):
        if sheets_url:
            with st.spinner("Loading…"):
                df_gs, err = try_load_google_sheet(sheets_url)
            if df_gs is not None:
                st.session_state.sheets_df = df_gs
                st.session_state.df_raw = None  # let Analyzer page pick it up fresh
                st.session_state.status_google_sheets = "Connected"
                st.success(f"Loaded {len(df_gs):,} rows × {len(df_gs.columns)} columns. Switch to the Analyzer page to use this data.")
                st.dataframe(df_gs.head(5), use_container_width=True)
            else:
                st.session_state.status_google_sheets = "Failed"
                st.error(f"Could not load sheet: {err}")
        else:
            st.warning("Enter a Google Sheets CSV export URL first.")

    with st.expander("Possible data pulls — Google Sheets"):
        for item in ["Creative performance table", "Campaign export table", "Testing tracker", "Weekly reporting sheet"]:
            st.markdown(f"- {item}")

    st.divider()

    # ── Meta Ads ─────────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>Meta Ads</div>", unsafe_allow_html=True)
    render_integration_card(
        title="Meta Ads",
        description="Direct source for spend, impressions, clicks, CTR, conversions, creative IDs, platform-level performance, and cost per result.",
        fields=[
            {"key": "access_token", "label": "Access Token", "password": True, "placeholder": "EAAxxxx…"},
            {"key": "ad_account_id", "label": "Ad Account ID", "placeholder": "act_123456789"},
        ],
        possible_pulls=["Creative-level spend", "CTR", "Impressions", "Clicks", "Conversions", "Cost per result", "Creative IDs", "Campaign and ad set metadata"],
        test_fn=test_meta_ads_connection,
        state_key="meta_ads",
    )

    # ── TikTok Ads ────────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>TikTok Ads</div>", unsafe_allow_html=True)
    render_integration_card(
        title="TikTok Ads",
        description="Direct source for short-form creative performance, thumbstop behavior, hold rates, CTR, spend, and conversion data.",
        fields=[
            {"key": "access_token", "label": "Access Token", "password": True, "placeholder": "Bearer xxxx…"},
            {"key": "advertiser_id", "label": "Advertiser ID", "placeholder": "1234567890123"},
        ],
        possible_pulls=["Thumbstop rate", "Hold rate", "CTR", "Spend", "Trial starts", "Paid starts", "Video-level engagement"],
        test_fn=test_tiktok_ads_connection,
        state_key="tiktok_ads",
    )

    # ── Airtable ──────────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>Airtable</div>", unsafe_allow_html=True)
    render_integration_card(
        title="Airtable",
        description="Useful for creative testing databases, hook libraries, concept tags, production status, test history, and learning logs.",
        fields=[
            {"key": "api_key", "label": "API Key", "password": True, "placeholder": "patXXXX…"},
            {"key": "base_id", "label": "Base ID", "placeholder": "appXXXXXXXX"},
        ],
        possible_pulls=["Creative concept database", "Hook library", "Test status", "Creative owner", "Production stage", "Learning tags"],
        test_fn=test_airtable_connection,
        state_key="airtable",
    )

    # ── Notion ────────────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>Notion</div>", unsafe_allow_html=True)
    render_integration_card(
        title="Notion",
        description="Useful for creative briefs, experiment documentation, weekly learnings, and team-facing summaries.",
        fields=[
            {"key": "api_key", "label": "Integration Token", "password": True, "placeholder": "secret_xxxx…"},
        ],
        possible_pulls=["Creative briefs", "Weekly growth notes", "Test summaries", "Experiment writeups"],
        test_fn=test_notion_connection,
        state_key="notion",
    )

    # ── Slack ─────────────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>Slack</div>", unsafe_allow_html=True)
    render_integration_card(
        title="Slack",
        description="Useful for sending scale, cut, and test-next recommendations directly to the growth team.",
        fields=[
            {"key": "bot_token", "label": "Bot Token", "password": True, "placeholder": "xoxb-xxxx…"},
        ],
        possible_pulls=["Scale and cut alerts", "Weekly creative summaries", "Test recommendations", "Team notifications"],
        test_fn=test_slack_connection,
        state_key="slack",
    )

    # ── Google Drive ──────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>Google Drive</div>", unsafe_allow_html=True)
    render_integration_card(
        title="Google Drive",
        description="Useful for linking creative assets, raw videos, thumbnails, briefs, and campaign exports.",
        fields=[
            {"key": "api_key", "label": "API Key", "password": True, "placeholder": "AIzaXXXX…"},
        ],
        possible_pulls=["Creative files", "Video assets", "Thumbnails", "Briefs", "Exported reports"],
        test_fn=test_google_drive_connection,
        state_key="google_drive",
    )

    # ── Mixpanel ──────────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>Mixpanel</div>", unsafe_allow_html=True)
    render_integration_card(
        title="Mixpanel",
        description="Useful for product-side activation data, trial behavior, user cohorts, and trial-to-paid funnel analysis.",
        fields=[
            {"key": "api_secret", "label": "API Secret", "password": True, "placeholder": "xxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxx"},
        ],
        possible_pulls=["Activation events", "Trial usage behavior", "Product engagement by cohort", "Trial-to-paid funnel events"],
        test_fn=test_mixpanel_connection,
        state_key="mixpanel",
    )

    # ── Amplitude ─────────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>Amplitude</div>", unsafe_allow_html=True)
    render_integration_card(
        title="Amplitude",
        description="Useful for product analytics, retention behavior, activation events, and user journey analysis.",
        fields=[
            {"key": "api_key", "label": "API Key", "password": True, "placeholder": "xxxxxxxxxxxxxxxxxxxxxxxxxxxxxxxx"},
        ],
        possible_pulls=["User journey analysis", "Retention metrics", "Activation cohorts", "Funnel drop-off data"],
        test_fn=test_amplitude_connection,
        state_key="amplitude",
    )

    # ── HubSpot ───────────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>HubSpot</div>", unsafe_allow_html=True)
    render_integration_card(
        title="HubSpot",
        description="Useful for lead quality, lifecycle stage, CRM source attribution, and downstream conversion feedback.",
        fields=[
            {"key": "private_app_token", "label": "Private App Token", "password": True, "placeholder": "pat-na1-xxxx…"},
        ],
        possible_pulls=["Lead source quality", "Lifecycle stage", "CRM conversion feedback", "Deal source attribution"],
        test_fn=test_hubspot_connection,
        state_key="hubspot",
    )

    st.divider()

    # ── Production improvements ───────────────────────────────────────────────
    st.markdown("<div class='section-header'>Production improvements</div>", unsafe_allow_html=True)
    improvements = [
        "OAuth instead of manual API keys",
        "Encrypted secret storage",
        "Environment variables for server-side secrets",
        "Role-based permissions",
        "Scheduled data refresh",
        "Source-specific data mapping",
        "API rate limit handling",
        "Error logging",
        "Secret scanning before deployment",
        "Clear separation between frontend inputs and backend data pulls",
    ]
    cols = st.columns(2)
    for i, item in enumerate(improvements):
        with cols[i % 2]:
            st.markdown(f"<div class='insight-box'>✦ {item}</div>", unsafe_allow_html=True)


def page_methodology():
    st.markdown(
        """
        <div style="display:flex;align-items:center;gap:12px;margin-bottom:0.2rem;">
          <span style="font-size:2rem;">📐</span>
          <span style="font-size:1.7rem;font-weight:800;letter-spacing:-0.02em;">Methodology</span>
        </div>
        <p style="color:#8A9BC8;margin-top:0;margin-bottom:1rem;font-size:0.95rem;">
          How KPIs are calculated, how creatives are ranked, and why each decision label means what it means.
        </p>
        """,
        unsafe_allow_html=True,
    )
    st.divider()

    # ── KPI formulas ─────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>KPI Formulas</div>", unsafe_allow_html=True)
    kpis = [
        ("Cost per Paid Start", "Spend ÷ Paid Starts", "Primary efficiency metric. Lower = better."),
        ("Cost per Trial Start", "Spend ÷ Trial Starts", "Upper-funnel efficiency. Useful for trial-led products."),
        ("Trial → Paid CVR", "Paid Starts ÷ Trial Starts", "Conversion quality from trial to paying customer."),
        ("Paid Starts per $1k", "Paid Starts ÷ Spend × 1,000", "Volume efficiency — how many paid starts per thousand dollars spent."),
        ("Trial Starts per $1k", "Trial Starts ÷ Spend × 1,000", "Upper-funnel volume efficiency."),
        ("CTR Efficiency Score", "CTR ÷ Cost per Paid Start", "How much click-through rate you get per dollar of CPA."),
        ("Thumbstop Efficiency", "Thumbstop Rate ÷ Cost per Paid Start", "How much scroll-stopping power you get per dollar of CPA."),
        ("Hold Efficiency", "6s Hold Rate ÷ Cost per Paid Start", "How much 6-second retention you get per dollar of CPA."),
    ]
    for name, formula, note in kpis:
        st.markdown(
            f"<div style='background:#161B27;border:1px solid #2A3350;border-radius:8px;"
            f"padding:12px 18px;margin-bottom:8px;'>"
            f"<span style='color:#4F8EF7;font-weight:700;'>{name}</span>"
            f"<span style='color:#8A9BC8;font-size:0.82rem;margin-left:12px;'>{formula}</span>"
            f"<div style='color:#C5CFDF;font-size:0.88rem;margin-top:4px;'>{note}</div>"
            f"</div>",
            unsafe_allow_html=True,
        )

    st.divider()

    # ── Ranking logic ─────────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>Ranking Logic by Campaign Goal</div>", unsafe_allow_html=True)
    goals_info = [
        ("Paid Starts", "Sort by Paid Starts descending. Best for pure volume goals."),
        ("Trial Starts", "Sort by Trial Starts descending. Best for upper-funnel volume."),
        ("Efficient Paid Starts", f"Primary: Cost per Paid Start ascending. Secondary: Paid Starts descending. Creatives below the ${CPA_TARGET} CPA target are flagged as Scale Candidates."),
        ("Efficient Trial Starts", "Primary: Cost per Trial Start ascending. Secondary: Trial Starts descending."),
        ("Creative Engagement", "Weighted score: 40% Thumbstop Rate + 30% 6s Hold Rate + 30% CTR. Each metric is min-max normalised before scoring. Best for hook and mid-funnel creative testing."),
        ("Full Funnel Quality", "Weighted normalised score: 30% Cost per Paid Start (lower is better) + 25% Paid Starts + 20% Trial→Paid CVR + 15% CTR + 10% Thumbstop Rate. Best for holistic creative reviews."),
    ]
    for goal_name, logic in goals_info:
        st.markdown(
            f"<div style='background:#161B27;border:1px solid #2A3350;border-radius:8px;"
            f"padding:12px 18px;margin-bottom:8px;'>"
            f"<span style='color:#4F8EF7;font-weight:700;'>{goal_name}</span>"
            f"<div style='color:#C5CFDF;font-size:0.88rem;margin-top:4px;'>{logic}</div>"
            f"</div>",
            unsafe_allow_html=True,
        )

    st.divider()

    # ── Decision labels ───────────────────────────────────────────────────────
    st.markdown("<div class='section-header'>Decision Labels</div>", unsafe_allow_html=True)
    label_info = [
        ("Scale", "#34d399", "#1a4731", f"CPA is below ${CPA_TARGET} target AND Paid Starts are above median. This creative is working — invest more."),
        ("Keep Testing", "#fbbf24", "#3b3515", "Engagement metrics are above median but Paid Starts are below median. The hook is working; the funnel needs attention."),
        ("Fix Funnel", "#fb923c", "#3b2200", "Trial Starts are above median but Trial→Paid CVR is below median. Traffic is coming through; conversion is the problem."),
        ("Cut", "#f87171", "#3b0a0a", f"CPA is above ${CPA_TARGET} AND engagement is below median. No signal worth investing in. Retire and reallocate."),
        ("Review", "#f59e0b", "#2e2a10", "Does not clearly meet any of the above criteria. Needs manual review before a decision."),
    ]
    for label, fg, bg, explanation in label_info:
        st.markdown(
            f"<div style='background:{bg};border:1px solid {fg};border-radius:8px;"
            f"padding:12px 18px;margin-bottom:8px;'>"
            f"<span style='color:{fg};font-weight:700;font-size:1rem;'>{label}</span>"
            f"<div style='color:#C5CFDF;font-size:0.88rem;margin-top:4px;'>{explanation}</div>"
            f"</div>",
            unsafe_allow_html=True,
        )

    st.divider()

    # ── Why CPA matters for LoopNote ──────────────────────────────────────────
    st.markdown("<div class='section-header'>Why Cost per Paid Start matters most for LoopNote</div>", unsafe_allow_html=True)
    st.markdown(
        "<div class='insight-box'>"
        "LoopNote is a subscription product with a trial-to-paid model. The most important signal is whether an ad creative "
        "drives paying customers — not just clicks, impressions, or even trial sign-ups. A creative with a low thumbstop rate "
        "but a sub-$90 CPA is more valuable than a viral hook that never converts. Cost per Paid Start cuts through surface-level "
        "engagement and connects creative investment directly to revenue."
        "</div>",
        unsafe_allow_html=True,
    )
    st.markdown(
        "<div class='insight-box'>"
        "That said, CPA alone can mislead on small spend samples. Creatives with fewer than 5–10 paid starts should be treated "
        "as directional signals, not conclusive results. Always read CPA alongside Paid Starts volume."
        "</div>",
        unsafe_allow_html=True,
    )

    st.divider()

    # ── Adapting to other goals ───────────────────────────────────────────────
    st.markdown("<div class='section-header'>How the tool adapts to other campaign goals</div>", unsafe_allow_html=True)
    st.markdown(
        "<div class='insight-box'>"
        "The goal selector re-weights the ranking logic without changing the underlying data. A team focused on trial volume "
        "can rank by Trial Starts; a team testing creative hooks before scaling can use Creative Engagement scoring. "
        "Full Funnel Quality is best for quarterly reviews where all five signals matter equally. Switching goals never alters "
        "the raw KPIs — only the sort order and score changes."
        "</div>",
        unsafe_allow_html=True,
    )

    st.divider()

    # ── Engagement vs conversion ──────────────────────────────────────────────
    st.markdown("<div class='section-header'>Why engagement metrics should be read alongside conversion metrics</div>", unsafe_allow_html=True)
    st.markdown(
        "<div class='insight-box'>"
        "Thumbstop Rate and 6s Hold Rate measure whether a creative earns attention in the first moments of a scroll. "
        "CTR measures whether it earns a click. But attention and clicks do not guarantee paid subscriptions. A creative "
        "can stop the scroll and still fail to convert if the offer is weak, the landing page creates friction, or the "
        "product doesn't match the promise in the ad."
        "</div>",
        unsafe_allow_html=True,
    )
    st.markdown(
        "<div class='insight-box'>"
        "The most useful pattern to look for is a creative with high engagement but high CPA — this is a funnel problem, "
        "not a creative problem. Conversely, a creative with low engagement but low CPA is doing something unusual: "
        "it may be self-selecting a high-intent audience. Both patterns are worth investigating before making scale or cut decisions."
        "</div>",
        unsafe_allow_html=True,
    )


# ══════════════════════════════════════════════════════════════════════════════
# PAGE: WEEK-OVER-WEEK COMPARE
# ══════════════════════════════════════════════════════════════════════════════

def page_compare():
    st.markdown(
        """
        <div style="display:flex;align-items:center;gap:12px;margin-bottom:0.2rem;">
          <span style="font-size:2rem;">📅</span>
          <span style="font-size:1.7rem;font-weight:800;letter-spacing:-0.02em;">Week-over-Week Comparison</span>
        </div>
        <p style="color:#8A9BC8;margin-top:0;margin-bottom:1rem;font-size:0.95rem;">
          Upload two data exports to compare creative performance across periods.
          Match is done on Creative ID — creatives that appear in only one period show blanks for the other.
        </p>
        """,
        unsafe_allow_html=True,
    )
    st.divider()

    col1, col2 = st.columns(2)
    with col1:
        st.markdown("**Current Period**")
        file_a = st.file_uploader(
            "Current period CSV / Excel",
            type=["csv", "xlsx", "xls"],
            key="wow_file_a",
        )
    with col2:
        st.markdown("**Previous Period**")
        file_b = st.file_uploader(
            "Previous period CSV / Excel",
            type=["csv", "xlsx", "xls"],
            key="wow_file_b",
        )

    if not file_a or not file_b:
        st.info("Upload both files above to generate the comparison.")
        return

    try:
        df_a, _ = load_and_clean(file_a)
        df_b, _ = load_and_clean(file_b)
    except Exception as e:
        st.error(f"Could not load files: {e}")
        return

    if "creative_id" not in df_a.columns or "creative_id" not in df_b.columns:
        st.error("Both files must have a Creative ID column.")
        return

    df_a, _ = calculate_kpis(df_a)
    df_b, _ = calculate_kpis(df_b)

    metrics = ["spend", "cpa", "paid_starts", "ctr", "thumbstop_rate", "trial_starts"]
    avail = [m for m in metrics if m in df_a.columns and m in df_b.columns]
    if not avail:
        st.warning("No common metric columns found between the two files.")
        return

    merged = df_a[["creative_id"] + avail].merge(
        df_b[["creative_id"] + avail],
        on="creative_id",
        suffixes=("_curr", "_prev"),
        how="outer",
    )

    metric_labels = {
        "spend": "Spend ($)", "cpa": "CPA ($)", "paid_starts": "Paid Starts",
        "ctr": "CTR (%)", "thumbstop_rate": "Thumbstop (%)", "trial_starts": "Trial Starts",
    }
    lower_is_better = {"cpa", "spend"}

    def _fmt(m, v):
        if pd.isna(v): return "—"
        if m in ("spend", "cpa"): return f"${v:,.0f}"
        if m in ("ctr", "thumbstop_rate"): return f"{v:.1f}%"
        return f"{int(v):,}"

    display_rows = []
    for _, row in merged.iterrows():
        r: dict = {"Creative": row["creative_id"]}
        for m in avail:
            curr = row.get(f"{m}_curr")
            prev = row.get(f"{m}_prev")
            lbl = metric_labels.get(m, m)
            r[f"{lbl} (Now)"] = _fmt(m, curr)
            r[f"{lbl} (Prev)"] = _fmt(m, prev)
            if pd.notna(curr) and pd.notna(prev) and prev != 0:
                pct = (curr - prev) / prev * 100
                arrow = "↑" if pct > 0 else "↓"
                r[f"{lbl} Δ"] = f"{arrow} {abs(pct):.1f}%"
            else:
                r[f"{lbl} Δ"] = "—"
        display_rows.append(r)

    st.markdown("<div class='section-header'>Performance Delta by Creative</div>", unsafe_allow_html=True)
    st.dataframe(pd.DataFrame(display_rows), use_container_width=True)

    # ── CPA comparison chart ───────────────────────────────────────────────
    if "cpa" in avail:
        cpa_data = merged[["creative_id", "cpa_curr", "cpa_prev"]].dropna()
        if not cpa_data.empty:
            cpa_melt = cpa_data.melt(
                id_vars="creative_id",
                value_vars=["cpa_curr", "cpa_prev"],
                var_name="Period", value_name="CPA ($)",
            )
            cpa_melt["Period"] = cpa_melt["Period"].map({"cpa_curr": "Current", "cpa_prev": "Previous"})
            fig_cpa = px.bar(
                cpa_melt, x="creative_id", y="CPA ($)", color="Period", barmode="group",
                color_discrete_map={"Current": "#4F8EF7", "Previous": "#8A9BC8"},
                labels={"creative_id": "Creative"},
                title="CPA: Current vs Previous Period",
            )
            fig_cpa.update_layout(
                paper_bgcolor="#0E1117", plot_bgcolor="#0E1117", font_color="#FAFAFA",
                xaxis=dict(gridcolor="#1E2A45", tickangle=-30),
                yaxis=dict(gridcolor="#1E2A45"),
            )
            st.plotly_chart(fig_cpa, use_container_width=True)

    # ── Paid Starts comparison chart ───────────────────────────────────────
    if "paid_starts" in avail:
        ps_data = merged[["creative_id", "paid_starts_curr", "paid_starts_prev"]].dropna()
        if not ps_data.empty:
            ps_melt = ps_data.melt(
                id_vars="creative_id",
                value_vars=["paid_starts_curr", "paid_starts_prev"],
                var_name="Period", value_name="Paid Starts",
            )
            ps_melt["Period"] = ps_melt["Period"].map({"paid_starts_curr": "Current", "paid_starts_prev": "Previous"})
            fig_ps = px.bar(
                ps_melt, x="creative_id", y="Paid Starts", color="Period", barmode="group",
                color_discrete_map={"Current": "#34d399", "Previous": "#8A9BC8"},
                labels={"creative_id": "Creative"},
                title="Paid Starts: Current vs Previous Period",
            )
            fig_ps.update_layout(
                paper_bgcolor="#0E1117", plot_bgcolor="#0E1117", font_color="#FAFAFA",
                xaxis=dict(gridcolor="#1E2A45", tickangle=-30),
                yaxis=dict(gridcolor="#1E2A45"),
            )
            st.plotly_chart(fig_ps, use_container_width=True)

    st.download_button(
        "⬇ Download Comparison CSV",
        data=pd.DataFrame(display_rows).to_csv(index=False).encode("utf-8"),
        file_name="wow_comparison.csv",
        mime="text/csv",
        key="dl_wow",
    )


# ══════════════════════════════════════════════════════════════════════════════
# APP LAYOUT
# ══════════════════════════════════════════════════════════════════════════════

# ── Initialise session state ─────────────────────────────────────────────────
if "df_raw" not in st.session_state:
    st.session_state.df_raw = None
if "sheets_df" not in st.session_state:
    st.session_state.sheets_df = None
if "loaded_goal" not in st.session_state:
    st.session_state.loaded_goal = None
if "uploader_key" not in st.session_state:
    st.session_state.uploader_key = 0
if "pending_load_id" not in st.session_state:
    st.session_state.pending_load_id = None
if "pending_delete_id" not in st.session_state:
    st.session_state.pending_delete_id = None

# ── Process pending load / delete (must happen before sidebar renders) ────────
if st.session_state.pending_load_id:
    _save = load_save(st.session_state.pending_load_id)
    if _save:
        try:
            _df, _warns = load_and_clean(StringIO(_save["csv_data"]))
            st.session_state.df_raw = _df
            st.session_state.warnings = _warns
            st.session_state.loaded_goal = _save.get("goal")
            if _save.get("cpa_target") is not None:
                st.session_state["cpa_target_input"] = float(_save["cpa_target"])
        except Exception as _load_err:
            st.error(f"Could not restore saved analysis: {_load_err}")
    st.session_state.pending_load_id = None

if st.session_state.pending_delete_id:
    delete_save(st.session_state.pending_delete_id)
    st.session_state.pending_delete_id = None

# ── Sidebar: navigation + advanced settings ───────────────────────────────────
with st.sidebar:
    st.markdown(
        "<div style='font-size:1.1rem;font-weight:800;letter-spacing:-0.01em;margin-bottom:4px;'>📊 Creative Analyzer</div>",
        unsafe_allow_html=True,
    )
    st.caption("Paid social creative intelligence")
    st.divider()
    page = st.radio(
        "Navigation",
        ["Analyzer", "Compare", "Integrations", "Methodology"],
        label_visibility="collapsed",
    )
    st.divider()
    st.caption("Use the Integrations page to connect data sources, or load the built-in LoopNote sample on the Analyzer page.")
    st.divider()

    # ── Saved analyses ────────────────────────────────────────────────────────
    st.markdown(
        "<div style='font-size:0.82rem;font-weight:700;color:#4F8EF7;text-transform:uppercase;"
        "letter-spacing:0.08em;margin-bottom:6px;'>💾 Saved Analyses</div>",
        unsafe_allow_html=True,
    )
    _all_saves = list_saves()
    if not _all_saves:
        st.caption("No saves yet. Run an analysis and save it from the Analyzer page.")
    else:
        for _s in _all_saves:
            _sid = _s["id"]
            _label = _html.escape(_s["name"])
            _ts = _html.escape(_s.get("saved_at", ""))
            _goal_tag = _html.escape(_s.get("goal", ""))
            st.markdown(
                f"<div style='background:#161B27;border:1px solid #2A3350;border-radius:8px;"
                f"padding:8px 12px;margin-bottom:6px;'>"
                f"<div style='font-size:0.88rem;font-weight:600;color:#FAFAFA;'>{_label}</div>"
                f"<div style='font-size:0.75rem;color:#8A9BC8;margin-top:2px;'>{_ts}"
                f"{' · ' + _goal_tag if _goal_tag else ''}</div>"
                f"</div>",
                unsafe_allow_html=True,
            )
            _load_col, _del_col = st.columns([3, 1])
            with _load_col:
                if st.button("Load", key=f"load_{_sid}", use_container_width=True):
                    st.session_state.pending_load_id = _sid
                    st.rerun()
            with _del_col:
                if st.button("✕", key=f"del_{_sid}", use_container_width=True, help="Delete this save"):
                    st.session_state.pending_delete_id = _sid
                    st.rerun()
    st.divider()
    with st.expander("⚙️ Advanced settings", expanded=False):
        st.markdown("**CPA Target**")
        cpa_target = st.number_input(
            "Target Cost per Paid Start ($)",
            min_value=1.0,
            max_value=10000.0,
            value=float(CPA_TARGET),
            step=5.0,
            key="cpa_target_input",
            help="Creatives below this CPA are flagged as Scale candidates. Affects decision labels and chart reference lines.",
        )

        st.markdown("---")
        st.markdown("**Benchmark Targets**")
        st.caption("Enter platform averages so the Benchmark Comparison expander can show deltas. Leave 0 to skip.")
        _bench_platforms = ["Meta", "TikTok", "YouTube", "Google"]
        _benchmarks: dict[str, float] = {}
        for _bp in _bench_platforms:
            bc1, bc2, bc3 = st.columns(3)
            with bc1:
                st.caption(_bp)
            with bc2:
                _bv_ctr = st.number_input(f"CTR % ({_bp})", min_value=0.0, value=0.0, step=0.1, key=f"bench_{_bp}_ctr", label_visibility="collapsed")
                if _bv_ctr > 0:
                    _benchmarks[f"{_bp}_ctr"] = _bv_ctr
            with bc3:
                _bv_cpa = st.number_input(f"CPA $ ({_bp})", min_value=0.0, value=0.0, step=5.0, key=f"bench_{_bp}_cpa", label_visibility="collapsed")
                if _bv_cpa > 0:
                    _benchmarks[f"{_bp}_cpa"] = _bv_cpa
        st.caption("Columns: Platform · CTR (%) · CPA ($)")
        st.session_state["_benchmarks"] = _benchmarks

        st.markdown("---")
        st.markdown("**Creative Engagement weights**")
        st.caption("Set the relative importance of each signal. Weights are auto-normalised so the total always equals 100%.")

        ce_thumbstop = st.slider("Thumbstop Rate", 0, 100, CE_WEIGHTS_DEFAULT["thumbstop_rate"], key="ce_thumbstop")
        ce_hold = st.slider("6s Hold Rate", 0, 100, CE_WEIGHTS_DEFAULT["hold_6s"], key="ce_hold")
        ce_ctr = st.slider("CTR", 0, 100, CE_WEIGHTS_DEFAULT["ctr"], key="ce_ctr")
        ce_total = ce_thumbstop + ce_hold + ce_ctr
        if ce_total == 0:
            st.error("At least one weight must be greater than 0.")
            ce_weights = CE_WEIGHTS_DEFAULT
        else:
            ce_weights = {
                "thumbstop_rate": ce_thumbstop,
                "hold_6s": ce_hold,
                "ctr": ce_ctr,
            }
            norm_thumb = round(ce_thumbstop / ce_total * 100)
            norm_hold = round(ce_hold / ce_total * 100)
            norm_ctr = 100 - norm_thumb - norm_hold
            st.caption(f"Normalised → Thumbstop {norm_thumb}% · Hold {norm_hold}% · CTR {norm_ctr}%")

        st.markdown("---")
        st.markdown("**Full Funnel Quality weights**")
        st.caption("Set the relative importance of each metric. Weights are auto-normalised so the total always equals 100%.")

        ffq_cpa = st.slider("CPA (lower is better)", 0, 100, FFQ_WEIGHTS_DEFAULT["cpa"], key="ffq_cpa")
        ffq_paid = st.slider("Paid Starts", 0, 100, FFQ_WEIGHTS_DEFAULT["paid_starts"], key="ffq_paid")
        ffq_cvr = st.slider("Trial→Paid CVR", 0, 100, FFQ_WEIGHTS_DEFAULT["trial_to_paid_cvr"], key="ffq_cvr")
        ffq_ctr = st.slider("CTR", 0, 100, FFQ_WEIGHTS_DEFAULT["ctr"], key="ffq_ctr")
        ffq_thumbstop = st.slider("Thumbstop Rate", 0, 100, FFQ_WEIGHTS_DEFAULT["thumbstop_rate"], key="ffq_thumbstop")
        ffq_total = ffq_cpa + ffq_paid + ffq_cvr + ffq_ctr + ffq_thumbstop
        if ffq_total == 0:
            st.error("At least one weight must be greater than 0.")
            ffq_weights = FFQ_WEIGHTS_DEFAULT
        else:
            ffq_weights = {
                "cpa": ffq_cpa,
                "paid_starts": ffq_paid,
                "trial_to_paid_cvr": ffq_cvr,
                "ctr": ffq_ctr,
                "thumbstop_rate": ffq_thumbstop,
            }
            n_cpa = round(ffq_cpa / ffq_total * 100)
            n_paid = round(ffq_paid / ffq_total * 100)
            n_cvr = round(ffq_cvr / ffq_total * 100)
            n_ctr = round(ffq_ctr / ffq_total * 100)
            n_thumb = 100 - n_cpa - n_paid - n_cvr - n_ctr
            st.caption(f"Normalised → CPA {n_cpa}% · Paid {n_paid}% · CVR {n_cvr}% · CTR {n_ctr}% · Thumbstop {n_thumb}%")

# ── Route to page ─────────────────────────────────────────────────────────────
if page == "Analyzer":
    page_analyzer(cpa_target=cpa_target, ce_weights=ce_weights, ffq_weights=ffq_weights)
elif page == "Compare":
    page_compare()
elif page == "Integrations":
    page_integrations()
elif page == "Methodology":
    page_methodology()
